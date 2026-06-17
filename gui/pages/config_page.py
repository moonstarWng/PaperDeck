"""
gui/pages/config_page.py — 输入文件页：PDF + 模板 + 图片目录 + 模板检测/提取 + 一键生成。
"""
import customtkinter as ctk
import sys, os, threading, time, json
sys.path.insert(0, os.path.join(os.path.dirname(__file__), '..', '..'))
from gui.widgets.file_picker import FilePicker
from gui.logger import log_step


class ConfigPage(ctk.CTkFrame):
    """Tab 1: 论文 PDF、模板 PPTX、图片目录。"""

    def __init__(self, master, shared, app):
        super().__init__(master)
        self.pack(fill="both", expand=True)
        self.shared = shared
        self.app = app

        # ═══ 论文 PDF ═══
        ctk.CTkLabel(self, text="论文 PDF", font=ctk.CTkFont(size=16, weight="bold")).pack(anchor="w", pady=(10, 5))
        self.pdf_picker = FilePicker(self, "论文PDF:", [("PDF", "*.pdf"), ("All", "*.*")],
                                     on_change=lambda p: self._validate_pdf(p))
        self.pdf_picker.pack(fill="x", padx=10, pady=5)
        self.pdf_status = ctk.CTkLabel(self, text="未选择", text_color="gray")
        self.pdf_status.pack(anchor="w", padx=20)

        # ═══ 模板 PPTX + 检测 ═══
        ctk.CTkLabel(self, text="模板 PPTX", font=ctk.CTkFont(size=16, weight="bold")).pack(anchor="w", pady=(15, 5))
        tmpl_frame = ctk.CTkFrame(self)
        tmpl_frame.pack(fill="x", padx=10, pady=5)
        self.tmpl_picker = FilePicker(tmpl_frame, "模板:", [("PPTX", "*.pptx"), ("All", "*.*")],
                                      on_change=lambda p: self._validate_tmpl(p))
        self.tmpl_picker.pack(side="left", fill="x", expand=True)
        self.detect_btn = ctk.CTkButton(tmpl_frame, text="检测", width=60, command=self._detect_template)
        self.detect_btn.pack(side="right", padx=5)
        self.tmpl_status = ctk.CTkLabel(self, text="未选择", text_color="gray")
        self.tmpl_status.pack(anchor="w", padx=20)

        # ── 提取模板按钮（检测到完整 PPTX 后显示）──
        self.extract_frame = ctk.CTkFrame(self)
        self.extract_btn = ctk.CTkButton(
            self.extract_frame, text="提取模板骨架", width=130,
            command=self._run_extract_template)
        self.extract_btn.pack(side="left", padx=(10, 5))
        self.extract_progress = ctk.CTkLabel(
            self.extract_frame, text="", text_color="gray")
        self.extract_progress.pack(side="left", padx=5)
        self.preview_btn = ctk.CTkButton(
            self.extract_frame, text="预览生成文件", width=100, fg_color="#2E7D32",
            command=self._preview_template)
        # preview_btn 默认不显示，提取成功后 pack

        # ── 模板提取模式 ──
        mode_frame = ctk.CTkFrame(self)
        mode_frame.pack(fill="x", padx=10, pady=(5, 0))
        ctk.CTkLabel(mode_frame, text="提取模式:", width=70).pack(side="left", padx=(5, 0))
        self.extract_mode = ctk.StringVar(value="pptm")
        self.mode_pptm_btn = ctk.CTkRadioButton(
            mode_frame, text="PPT Master (推荐)", variable=self.extract_mode, value="pptm",
            command=self._on_mode_change)
        self.mode_pptm_btn.pack(side="left", padx=5)
        self.mode_rule_btn = ctk.CTkRadioButton(
            mode_frame, text="规则 (快速)", variable=self.extract_mode, value="rule",
            command=self._on_mode_change)
        self.mode_rule_btn.pack(side="left", padx=5)
        self.mode_llm_btn = ctk.CTkRadioButton(
            mode_frame, text="LLM (自适应)", variable=self.extract_mode, value="llm",
            command=self._on_mode_change)
        self.mode_llm_btn.pack(side="left", padx=5)
        self.mode_hint = ctk.CTkLabel(mode_frame, text="", text_color="gray")
        self.mode_hint.pack(side="left", padx=10)

        # ═══ 图片目录 ═══
        ctk.CTkLabel(self, text="论文图片目录 (figs/)", font=ctk.CTkFont(size=16, weight="bold")).pack(anchor="w", pady=(15, 5))
        self.figs_picker = FilePicker(self, "图片:", directory=True,
                                      on_change=lambda p: self._validate_figs(p))
        self.figs_picker.pack(fill="x", padx=10, pady=5)
        self.figs_status = ctk.CTkLabel(self, text="(可选)", text_color="gray")
        self.figs_status.pack(anchor="w", padx=20)

        # ═══ 一键生成 + 精细控制 ═══
        btn_frame = ctk.CTkFrame(self)
        btn_frame.pack(fill="x", padx=10, pady=(15, 5))

        self.one_click_btn = ctk.CTkButton(
            btn_frame, text="一键生成 PPT", height=44, width=200,
            fg_color="#2E7D32", font=ctk.CTkFont(size=16, weight="bold"),
            command=self._one_click_generate)
        self.one_click_btn.pack(side="left", padx=5)

        self.one_click_stop_btn = ctk.CTkButton(
            btn_frame, text="停止", height=44, width=60,
            fg_color="#C0392B", command=self._stop_one_click)
        # 默认隐藏，生成时显示

        self.one_click_status = ctk.CTkLabel(btn_frame, text="", text_color="gray")
        self.one_click_status.pack(side="left", padx=10)

        self.expert_cb = ctk.CTkCheckBox(
            btn_frame, text="精细控制", width=90,
            command=self._toggle_expert)
        self.expert_cb.pack(side="right", padx=10)

        # 原下一步按钮（专家模式下显示）
        self.next_btn = ctk.CTkButton(self, text="下一步: 生成大纲 →", height=36, command=self._go_next)
        # 默认隐藏，勾选精细控制后显示

    # ── 文件合法性实时检测 ──
    def _validate_pdf(self, path):
        self.shared['pdf_path'] = path if path and os.path.exists(path) else ''
        if not path:
            self.pdf_status.configure(text="未选择", text_color="gray")
        elif os.path.exists(path):
            self.pdf_status.configure(text="✓ 文件有效", text_color="green")
        else:
            self.pdf_status.configure(text="✗ 文件不存在", text_color="red")
        self._persist_now()

    def _validate_tmpl(self, path):
        self.shared['template_path'] = path if path and os.path.exists(path) else ''
        if not path:
            self.tmpl_status.configure(text="未选择", text_color="gray")
        elif os.path.exists(path):
            self.tmpl_status.configure(text="✓ 文件有效（点击「检测」判断是否为模板骨架）", text_color="green")
        else:
            self.tmpl_status.configure(text="✗ 文件不存在", text_color="red")
        self._persist_now()

    def _validate_figs(self, path):
        self.shared['figs_dir'] = path if path and os.path.exists(path) else ''
        if not path:
            self.figs_status.configure(text="(可选)", text_color="gray")
        elif os.path.exists(path):
            self.figs_status.configure(text="✓ 目录有效", text_color="green")
        else:
            self.figs_status.configure(text="✗ 目录不存在", text_color="red")
        self._persist_now()

    def _persist_now(self):
        try:
            from gui.persistence import save_from_shared
            save_from_shared(self.shared)
        except Exception:
            pass

    def _validate_files(self):
        self._validate_pdf(self.pdf_picker.get_path())
        self._validate_tmpl(self.tmpl_picker.get_path())
        self._validate_figs(self.figs_picker.get_path())

    # ── 模式切换 ──
    def _on_mode_change(self):
        mode = self.extract_mode.get()
        if mode == "pptm":
            self.mode_hint.configure(text="无需 API，10-30s，自适应模板设计 (PPT Master)")
        elif mode == "llm":
            self.mode_hint.configure(text="需 API 连接，约 10-30s，自适应不同模板")
        else:
            self.mode_hint.configure(text="无需 API，瞬间完成，适合标准模板")
        self.extract_frame.pack_forget()
        self.preview_btn.pack_forget()
        self.shared['is_template'] = False
        self.tmpl_status.configure(text="模式已切换，请重新检测模板", text_color="gray")

    # ── 模板检测 ──
    def _detect_template(self):
        """检测 PPTX 是否为模板骨架。"""
        self._validate_files()
        log_step('config', f'模板检测开始 ({self.extract_mode.get()} 模式)')
        path = self.tmpl_picker.get_path()
        if not path:
            self.tmpl_status.configure(text="请先选择模板文件", text_color="red")
            return
        if not os.path.exists(path):
            self.tmpl_status.configure(text="✗ 模板文件不存在", text_color="red")
            return
        self.shared['template_path'] = path

        if self.extract_mode.get() == "llm":
            self._detect_template_llm(path)
        elif self.extract_mode.get() == "pptm":
            self._detect_template_pptm(path)
        else:
            self._detect_template_rule(path)

    def _detect_template_rule(self, path):
        try:
            from pptx import Presentation
            prs = Presentation(path)
            from scripts.make_template import classify
            t_count, c_count, placeholder_count = 0, 0, 0
            for slide in prs.slides:
                cat = classify(slide)
                if cat.startswith('SECTION_') or cat in ('COVER', 'TOC', 'THANKS', 'SUMMARY'):
                    t_count += 1
                elif cat in ('CONTENT', 'CONTENT_FRAME', 'DATA', 'DISCUSSION'):
                    c_count += 1
                    all_text = ' '.join(sh.text_frame.text for sh in slide.shapes if sh.has_text_frame)
                    if '此处填充' in all_text or '章节标题' in all_text or '论文标题' in all_text:
                        placeholder_count += 1
            if c_count > 0 and placeholder_count / c_count > 0.5:
                is_tpl = True
                ratio = 0.0
            else:
                ratio = c_count / max(len(prs.slides), 1)
                is_tpl = ratio < 0.2
            self.shared['is_template'] = is_tpl
            self._update_extract_ui(is_tpl, ratio)
        except Exception as e:
            self.tmpl_status.configure(text=f"检测失败: {e}", text_color="red")
            self.extract_frame.pack_forget()

    def _update_extract_ui(self, is_tpl, ratio=None):
        self.preview_btn.pack_forget()
        if is_tpl:
            self.tmpl_status.configure(
                text=f"✓ 已是模板骨架，可直接下一步",
                text_color="green")
            self.extract_frame.pack_forget()
            # 检查是否需要显示预览按钮（模板路径是否在 process/ 下）
            self._maybe_show_preview()
        else:
            mode = self.extract_mode.get()
            mode_label = {"pptm": "PPT Master", "llm": "LLM 自适应"}.get(mode, "规则")
            self.tmpl_status.configure(
                text=f"⚠ 检测到完整PPTX (内容页{ratio:.0%})，请点击「提取模板骨架」",
                text_color="orange")
            self.extract_frame.pack(fill="x", padx=10, pady=(5, 0), before=self.mode_rule_btn.master)
            self.extract_btn.configure(text=f"提取模板骨架 ({mode_label})", state="normal")
            self.extract_progress.configure(text="")

    def _maybe_show_preview(self):
        """如果模板路径在 process/ 目录下（说明是生成的文件），显示预览按钮。"""
        tmpl = self.shared.get('template_path', '')
        if tmpl and os.path.exists(tmpl) and 'process' in tmpl.replace('\\', '/'):
            self.extract_frame.pack(fill="x", padx=10, pady=(5, 0), before=self.mode_rule_btn.master)
            self.preview_btn.pack(side="left", padx=5)

    def _detect_template_llm(self, path):
        try:
            from scripts.llm_template_extract import load_template_cache
            cache = load_template_cache(path)
            if cache:
                cats = {s['index']: s['classification'] for s in cache['slides']}
                c_count = sum(1 for v in cats.values() if v == 'CONTENT')
                ratio = c_count / max(len(cats), 1)
                is_tpl = ratio < 0.2
                self.shared['is_template'] = is_tpl
                self._update_extract_ui(is_tpl, ratio)
                self.tmpl_status.configure(text=f"✓ 从缓存识别", text_color="green")
                return
        except Exception:
            pass

        try:
            from pptx import Presentation
            prs = Presentation(path)
            from scripts.make_template import classify
            c_count, ph_count = 0, 0
            for slide in prs.slides:
                cat = classify(slide)
                if cat in ('CONTENT', 'CONTENT_FRAME', 'DATA', 'DISCUSSION'):
                    c_count += 1
                    all_text = ' '.join(sh.text_frame.text for sh in slide.shapes if sh.has_text_frame)
                    if '此处填充' in all_text or '章节标题' in all_text or '论文标题' in all_text:
                        ph_count += 1
            if c_count > 0 and ph_count / c_count > 0.5:
                self.shared['is_template'] = True
                self._update_extract_ui(True, 0.0)
                self.tmpl_status.configure(text="✓ 已是模板骨架（检测到占位符），无需处理", text_color="green")
                from scripts.llm_template_extract import analyze_and_cache
                threading.Thread(target=lambda: analyze_and_cache(path), daemon=True).start()
                return
        except Exception:
            pass

        url = self.shared.get('api_base_url', '').strip()
        key = self.shared.get('api_key', '').strip()
        model = self.shared.get('api_model', '').strip()
        if not url or not key:
            self.tmpl_status.configure(text="LLM 模式需要填写 API 配置（点击右下角「AI配置」）", text_color="red")
            return
        self.tmpl_status.configure(text="LLM 分类中...", text_color="gray")
        threading.Thread(target=self._do_llm_detect, args=(path, url, key, model), daemon=True).start()

    def _do_llm_detect(self, path, url, key, model):
        try:
            from pptx import Presentation
            from scripts.llm_template_extract import (
                dump_all_slides, build_classification_prompt,
                parse_classification_response, CLASSIFICATION_SYSTEM,
            )
            import requests
            prs = Presentation(path)
            slides_data = dump_all_slides(prs)
            prompt = build_classification_prompt(slides_data)
            headers = {'Content-Type': 'application/json'}
            if key:
                headers['Authorization'] = f'Bearer {key}'
            resp = requests.post(
                f'{url.rstrip("/")}/chat/completions',
                headers=headers,
                json={'model': model, 'messages': [
                    {'role': 'system', 'content': CLASSIFICATION_SYSTEM},
                    {'role': 'user', 'content': prompt},
                ], 'temperature': 0.0},
                timeout=120,
            )
            resp.raise_for_status()
            content = resp.json()['choices'][0]['message']['content']
            classifications = parse_classification_response(content)
            t_count = sum(1 for v in classifications.values() if v != 'CONTENT')
            c_count = sum(1 for v in classifications.values() if v == 'CONTENT')
            total = len(classifications)
            ratio = c_count / max(total, 1)
            is_tpl = ratio < 0.2
            self.shared['is_template'] = is_tpl
            self.shared['llm_classifications'] = classifications
            cats = []
            for i in sorted(classifications.keys()):
                cats.append(f"p{i}={classifications[i]}")
            detail = ", ".join(cats)
            self.after(0, lambda: self._update_extract_ui_llm(is_tpl, ratio, detail))
        except Exception as e:
            self.after(0, lambda: self.tmpl_status.configure(
                text=f"LLM 检测失败: {str(e)[:120]}", text_color="red"))

    def _update_extract_ui_llm(self, is_tpl, ratio, detail):
        self._update_extract_ui(is_tpl, ratio)
        current = self.tmpl_status.cget("text")
        self.tmpl_status.configure(text=f"{current}  [{detail}]")

    # ── PPT Master 检测 ──
    def _detect_template_pptm(self, path):
        """PPT Master 模式：用 ppt-master 分析模板，判断是否已是骨架。"""
        try:
            from scripts.ppt_master_adapter import is_available, detect_is_template, _LAST_ERROR
            if not is_available():
                err = _LAST_ERROR or "ppt-master 目录未找到"
                self.tmpl_status.configure(
                    text=f"PPT Master 不可用: {err}",
                    text_color="orange")
                return

            self.tmpl_status.configure(text="PPT Master 分析中...", text_color="gray")
            is_tpl = detect_is_template(path)
            self.shared['is_template'] = is_tpl
            if is_tpl:
                self._update_extract_ui(True, 0.0)
                self.tmpl_status.configure(
                    text="✓ PPT Master: 已是模板骨架，可直接下一步",
                    text_color="green")
            else:
                from pptx import Presentation
                prs = Presentation(path)
                total = len(prs.slides)
                from scripts.make_template import classify
                c_count = sum(1 for s in prs.slides if classify(s) in ('CONTENT', 'CONTENT_FRAME', 'DATA', 'DISCUSSION'))
                ratio = c_count / max(total, 1)
                self._update_extract_ui(False, ratio)
                self.tmpl_status.configure(
                    text=f"PPT Master: 完整PPTX ({total}页, 内容{c_count}页)，请点击「提取模板骨架」",
                    text_color="orange")
        except Exception as e:
            self.tmpl_status.configure(text=f"PPT Master 检测失败: {e}", text_color="red")
            self.extract_frame.pack_forget()

    # ── 预览生成模板 ──
    def _preview_template(self):
        """用系统默认程序打开生成的模板 PPTX。"""
        path = self.shared.get('template_path', '')
        if path and os.path.exists(path):
            log_step('config', f'预览模板: {path}')
            try:
                os.startfile(path)
            except Exception as e:
                from tkinter import messagebox
                messagebox.showerror("无法打开", str(e))

    # ── 恢复共享数据 ──
    def restore_from_shared(self):
        if self.shared.get('pdf_path'):
            self.pdf_picker.set_path(self.shared['pdf_path'])
        if self.shared.get('template_path'):
            self.tmpl_picker.set_path(self.shared['template_path'])
        if self.shared.get('figs_dir'):
            self.figs_picker.set_path(self.shared['figs_dir'])
        if self.shared.get('extract_mode'):
            self.extract_mode.set(self.shared['extract_mode'])
            self._on_mode_change()
        self._validate_files()

    # ── 精细控制 ──
    def _toggle_expert(self):
        """切换专家模式：显示/隐藏大纲页和构建页。"""
        show = self.expert_cb.get()
        self.app.toggle_expert_mode(show)
        if show:
            self.next_btn.pack(pady=5)
            self.tmpl_status.configure(text="分步模式：可逐页调整模板、大纲和构建细节", text_color="gray")
        else:
            self.next_btn.pack_forget()
            self.tmpl_status.configure(text="", text_color="gray")

    # ── 一键生成 ──
    def _one_click_generate(self):
        """一键生成：第一击生成大纲，第二击保存 PPT。"""
        # 如果大纲已准备好，直接保存
        if self.shared.get('_one_click_ready'):
            self._save_after_generate()
            return

        self._validate_files()

        # 校验文件
        pdf_path = self.pdf_picker.get_path()
        tmpl_path = self.tmpl_picker.get_path()
        if not pdf_path or not os.path.exists(pdf_path):
            self.one_click_status.configure(text="请先选择论文 PDF", text_color="red")
            return
        if not tmpl_path or not os.path.exists(tmpl_path):
            self.one_click_status.configure(text="请先选择模板 PPTX", text_color="red")
            return

        # 检查 API
        key = self.shared.get('api_key', '').strip()
        if not key:
            self.one_click_status.configure(text="请先配置 AI API（右下角 AI配置）", text_color="red")
            return

        # 锁定 UI
        self._stop_requested = False
        self.one_click_btn.configure(state="disabled", text="生成中...")
        self.one_click_stop_btn.pack(side="left", padx=5, after=self.one_click_btn)
        self.expert_cb.configure(state="disabled")
        self.detect_btn.configure(state="disabled")
        self.next_btn.configure(state="disabled")
        self.one_click_status.configure(text="开始一键生成...", text_color="gray")
        self.shared['pdf_path'] = pdf_path
        self.shared['template_path'] = tmpl_path
        self.shared['figs_dir'] = self.figs_picker.get_path() or ''

        threading.Thread(target=self._do_one_click, args=(pdf_path, tmpl_path), daemon=True).start()

    def _stop_one_click(self):
        """停止一键生成。"""
        self._stop_requested = True
        # 重置 outline_page 状态，中断流水线
        ol = self.app.outline_page
        if ol._state != ol.STATE_READY:
            ol._state = ol.STATE_IDLE
        self.one_click_status.configure(text="正在停止...", text_color="orange")

    def _do_one_click(self, pdf_path, tmpl_path):
        """一键生成后台线程：复用 OutlinePage 的任务流水线。"""
        def progress(msg):
            self.one_click_status.configure(text=msg)

        def stopped():
            return getattr(self, '_stop_requested', False)

        try:
            ol = self.app.outline_page

            # 清空上次的旧数据，防止残留
            for key in ('slide_content_json', 'paper_text', 'paper_meta', 'paper_title', '_one_click_ready'):
                self.shared.pop(key, None)

            # ── Step 1: 读取论文 + 提取元数据（等同于自动点"读取论文"按钮）──
            if stopped(): return self._finish_one_click("已停止")
            progress("Step 1/2: 读取论文并提取信息...")
            try:
                # 重置章节配置为默认值
                ol._reset_sections()
                # 直接调用读取论文的实际逻辑（不走 UI 线程）
                ol._do_read_paper(pdf_path)
                # 等待元数据提取完成（_extract_metadata 启动独立线程，需等待它结束）
                progress("Step 1/2: 等待元数据...")
                waited = 0
                while ol._state != ol.STATE_READY and waited < 30:
                    if stopped(): return self._finish_one_click("已停止")
                    time.sleep(0.3)
                    waited += 0.3
                # 从 shared dict 读标题（元数据线程已写入，比读 tk widget 更可靠）
                paper_meta = self.shared.get('paper_meta', {})
                title = paper_meta.get('title_en', '') if paper_meta else ''
                if title:
                    self.shared['paper_title'] = title
                    progress(f"Step 1/2: {title[:30]}...")
                else:
                    progress(f"Step 1/2: 读取完成（标题待定）")
            except Exception as e:
                progress(f"Step 1/2: 读取失败 ({e})")

            # ── Step 2: 生成大纲（等同于自动点"生成大纲"按钮）──
            if stopped(): return self._finish_one_click("已停止")
            progress("Step 2/2: 生成 PPT 大纲（调用 LLM，约需30-90秒）...")
            try:
                self._run_outline_pipeline()
                progress("完成！点击「一键生成 PPT」保存文件")
            except Exception as e:
                progress(f"大纲生成失败: {e}")
                import traceback; traceback.print_exc()
                return self._finish_one_click("大纲生成失败")

            # 大纲生成成功后，按钮文字变更提示用户下一步操作
            self.one_click_btn.configure(text="点击保存 PPT", fg_color="#1565C0")
            self.shared['_one_click_ready'] = True

        except Exception as e:
            progress(f"失败: {str(e)[:80]}")
            import traceback; traceback.print_exc()
        finally:
            self._finish_one_click(None)

    def _run_outline_pipeline(self):
        """调用 OutlinePage 的任务流水线生成 slide-content.json。"""
        ol = self.app.outline_page

        if not ol._check_api():
            raise RuntimeError("API 未配置")

        # 调用现有的任务流水线（和用户手动点"生成大纲"完全一样）
        import time
        ol._generate_outline()

        # 等待流水线完成（最多 5 分钟）
        waited = 0
        while ol._state != ol.STATE_READY and waited < 300:
            if getattr(self, '_stop_requested', False):
                ol._state = ol.STATE_IDLE  # 重置状态
                raise InterruptedError("用户停止")
            time.sleep(0.5)
            waited += 0.5

        if ol._state != ol.STATE_READY:
            raise RuntimeError("大纲生成超时或失败")

    def _finish_one_click(self, msg):
        """清理一键生成状态。"""
        if msg:
            self.one_click_status.configure(text=msg, text_color="orange" if "败" in msg or "止" in msg else "gray")
        # 如果生成成功则保持蓝色按钮让用户点击保存
        if not self.shared.get('_one_click_ready'):
            self.one_click_btn.configure(state="normal", text="一键生成 PPT", fg_color="#2E7D32")
        else:
            self.one_click_btn.configure(state="normal")
        self.one_click_stop_btn.pack_forget()
        self.expert_cb.configure(state="normal")
        self.detect_btn.configure(state="normal")
        self.next_btn.configure(state="normal")

    def _save_after_generate(self):
        """大纲生成完成后，用户点按钮保存 PPT。"""
        from tkinter import filedialog
        pdf_path = self.shared.get('pdf_path', '')

        output_path = filedialog.asksaveasfilename(
            title="保存 PPT",
            defaultextension=".pptx",
            filetypes=[("PowerPoint", "*.pptx")],
            initialfile=os.path.splitext(os.path.basename(pdf_path))[0] + "_汇报.pptx",
        )
        if not output_path:
            return

        # 构建 PPT
        json_str = self.shared.get('slide_content_json', '')
        config = json.loads(json_str)
        config.setdefault('meta', {})
        config['meta']['template_path'] = self.shared.get('template_path', '')
        config['meta']['figs_dir'] = self.shared.get('figs_dir', '')
        config['meta']['output_path'] = output_path
        config['meta']['paper_meta'] = self.shared.get('paper_meta', {})

        import tempfile
        tmp = tempfile.NamedTemporaryFile(mode='w', suffix='.json', delete=False, encoding='utf-8')
        json.dump(config, tmp, ensure_ascii=False, indent=2)
        tmp_path = tmp.name
        tmp.close()

        from scripts.ppt_builder import build
        try:
            build(config, tmp_path)
            os.unlink(tmp_path)
            self.one_click_status.configure(text=f"完成！→ {os.path.basename(output_path)}", text_color="green")
            try:
                os.startfile(output_path)
            except Exception:
                pass
        except Exception as e:
            self.one_click_status.configure(text=f"构建失败: {e}", text_color="red")

        # 恢复按钮
        self.shared['_one_click_ready'] = False
        self.one_click_btn.configure(text="一键生成 PPT", fg_color="#2E7D32", command=self._one_click_generate)

    # ── 下一步（专家模式）──
    def _go_next(self):
        # 先校验文件（刷新状态标签）
        self._validate_files()

        errors = []
        if not self.pdf_picker.get_path():
            errors.append("论文 PDF 未选择")
        elif not self.pdf_picker.is_valid():
            errors.append("论文 PDF 文件不存在")
        if not self.tmpl_picker.get_path():
            errors.append("模板 PPTX 未选择")
        elif not self.tmpl_picker.is_valid():
            errors.append("模板 PPTX 文件不存在")
        figs = self.figs_picker.get_path()
        if figs and not self.figs_picker.is_valid():
            errors.append("图片目录不存在")
        key = self.shared.get('api_key', '').strip()
        if not key:
            errors.append("API Key 未填写（点击右下角「AI配置」）")

        if errors:
            from tkinter import messagebox
            messagebox.showerror("配置不完整", "\n".join(errors))
            return

        self.shared['pdf_path'] = self.pdf_picker.get_path()
        self.shared['template_path'] = self.tmpl_picker.get_path()
        self.shared['figs_dir'] = figs
        self.shared['extract_mode'] = self.extract_mode.get()

        # 模板未提取为骨架：宽松处理，允许继续
        if not self.shared.get('is_template', False):
            from tkinter import messagebox
            ok = messagebox.askyesno(
                "模板未处理",
                "当前模板可能尚未提取为骨架，\n"
                "直接使用完整 PPTX 可能导致样式不一致。\n\n"
                "是否仍然继续？\n"
                "（建议: 点「否」返回，先点击「提取模板骨架」）")
            if not ok:
                return

        self.app.switch_to_tab("大纲生成")

    # ── 提取模板 ──
    def _run_extract_template(self):
        log_step('config', f'开始提取模板 ({self.extract_mode.get()} 模式)')
        self.extract_btn.configure(state="disabled", text="提取中...")
        self.extract_progress.configure(text="准备中...", text_color="gray")
        self.preview_btn.pack_forget()
        self.next_btn.configure(state="disabled")
        mode = self.extract_mode.get()
        if mode == "llm":
            threading.Thread(target=self._do_extract_llm, daemon=True).start()
        elif mode == "pptm":
            threading.Thread(target=self._do_extract_pptm, daemon=True).start()
        else:
            threading.Thread(target=self._do_extract_rule, daemon=True).start()

    def _on_extract_done(self, dst):
        """提取成功后的 UI 更新（显示预览按钮）。"""
        self.extract_btn.configure(text="提取模板骨架", state="disabled")
        self.next_btn.configure(state="normal")
        if os.path.exists(dst):
            self.preview_btn.pack(side="left", padx=5)

    def _do_extract_rule(self):
        try:
            self.extract_progress.configure(text="提取中...")
            import sys, os as _os
            src = self.shared['template_path']
            proc_dir = _os.path.join(_os.path.dirname(_os.path.abspath(src)), 'process')
            _os.makedirs(proc_dir, exist_ok=True)
            dst = _os.path.join(proc_dir, 'template.pptx')
            sys.argv = ['make_template.py', src, dst]
            from scripts.make_template import main as make_tmpl_main
            make_tmpl_main()
            try:
                from scripts.template_extractor import extract as extract_design
                design = extract_design(dst)
                import json as _json
                design_path = dst.replace('.pptx', '_design.json')
                with open(design_path, 'w', encoding='utf-8') as f:
                    _json.dump(design, f, indent=2, ensure_ascii=False)
            except Exception:
                pass
            self.shared['template_path'] = dst
            self.shared['is_template'] = True
            self.tmpl_status.configure(
                text=f"✓ 模板骨架已提取 → {os.path.basename(dst)}", text_color="green")
            self.extract_progress.configure(text="完成", text_color="green")
            self._on_extract_done(dst)
        except Exception as e:
            self.extract_progress.configure(text=f"失败: {str(e)[:80]}", text_color="red")
            self.extract_btn.configure(state="normal")
            self.next_btn.configure(state="normal")

    def _do_extract_llm(self):
        def progress(msg):
            self.extract_progress.configure(text=msg, text_color="gray")

        try:
            import os as _os
            src = self.shared['template_path']
            proc_dir = _os.path.join(_os.path.dirname(_os.path.abspath(src)), 'process')
            _os.makedirs(proc_dir, exist_ok=True)
            dst = _os.path.join(proc_dir, 'template.pptx')
            url = self.shared.get('api_base_url', '').strip()
            key = self.shared.get('api_key', '').strip()
            model = self.shared.get('api_model', '').strip()

            from scripts.llm_template_extract import extract_template_llm
            result = extract_template_llm(src, dst, url, key, model, on_progress=progress)

            try:
                from scripts.template_extractor import extract as extract_design
                design = extract_design(dst)
                import json as _json
                design_path = dst.replace('.pptx', '_design.json')
                with open(design_path, 'w', encoding='utf-8') as f:
                    _json.dump(design, f, indent=2, ensure_ascii=False)
            except Exception:
                pass

            self.shared['template_path'] = dst
            self.shared['is_template'] = True
            self.shared['llm_classifications'] = result['classifications']
            self.tmpl_status.configure(
                text=f"✓ LLM 模板已提取 → {os.path.basename(dst)}", text_color="green")
            self.extract_progress.configure(text="完成", text_color="green")
            self.extract_frame.pack_forget()
            self._on_extract_done(dst)
        except Exception as e:
            self.extract_progress.configure(text=f"失败: {str(e)[:80]}", text_color="red")
            self.extract_btn.configure(state="normal")
            self.next_btn.configure(state="normal")

    def _do_extract_pptm(self):
        """PPT Master 模式：分析原始模板，提取增强设计 token。"""
        try:
            import os as _os
            src = self.shared['template_path']
            proc_dir = _os.path.join(_os.path.dirname(_os.path.abspath(src)), 'process')
            _os.makedirs(proc_dir, exist_ok=True)

            # Step 1: 分析原始模板
            self.extract_progress.configure(text="PPT Master 分析中...")
            from scripts.ppt_master_adapter import analyze_template, extract_design_tokens, save_analysis_cache
            analysis = analyze_template(src)

            # Step 2: 提取增强设计 token
            self.extract_progress.configure(text="提取设计 token...")
            design = extract_design_tokens(src, analysis)

            # Step 3: 保存到 process/
            import json as _json
            design_path = _os.path.join(proc_dir, 'template_design.json')
            with open(design_path, 'w', encoding='utf-8') as f:
                _json.dump(design, f, indent=2, ensure_ascii=False)
            cache_path = _os.path.join(proc_dir, 'analysis_pptm.json')
            save_analysis_cache(analysis, cache_path)

            # Step 4: 更新状态
            self.shared['template_path'] = src
            self.shared['is_template'] = True
            self.shared['pptmaster_analysis'] = analysis

            page_info = f"({analysis['slide_count']}页)"
            short_map = {'cover_candidate': '封面', 'toc_candidate': '目录',
                         'chapter_candidate': '章节', 'content_candidate': '内容',
                         'ending_candidate': '致谢'}
            for idx in sorted(analysis.get('page_types', {}).keys()):
                pt = analysis['page_types'][idx]
                short = short_map.get(pt, pt)
                page_info += f" p{idx}={short}"

            self.tmpl_status.configure(
                text=f"PPT Master 分析完成 {page_info}", text_color="green")
            self.extract_progress.configure(text="完成", text_color="green")
            self._on_extract_done(src)

        except Exception as e:
            import traceback
            traceback.print_exc()
            self.extract_progress.configure(
                text=f"PPT Master 失败，回退到规则模式...", text_color="orange")
            try:
                self._do_extract_rule()
            except Exception:
                self.extract_progress.configure(
                    text=f"失败: {str(e)[:80]}", text_color="red")
                self.extract_btn.configure(state="normal")
                self.next_btn.configure(state="normal")
