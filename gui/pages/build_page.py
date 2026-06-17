"""
gui/pages/build_page.py — 构建页：选择输出位置 + 验证 JSON + 生成 PPT + 进度日志。
"""
import customtkinter as ctk
import json, sys, os, threading, tempfile
from tkinter import messagebox, filedialog
from gui.logger import log_step, log_error

sys.path.insert(0, os.path.join(os.path.dirname(__file__), '..', '..'))


class BuildPage(ctk.CTkFrame):
    """Tab 3: 构建触发 + 输出选择 + 实时日志。"""

    def __init__(self, master, shared, app):
        super().__init__(master)
        self.pack(fill="both", expand=True)
        self.shared = shared
        self.app = app
        self.output_path = ""

        # ── 顶部按钮栏 ──
        btn_frame = ctk.CTkFrame(self)
        btn_frame.pack(fill="x", padx=10, pady=10)

        ctk.CTkButton(btn_frame, text="← 返回大纲", width=100, command=lambda: app.switch_to_tab("大纲生成")).pack(side="left", padx=5)
        self.build_btn = ctk.CTkButton(btn_frame, text="▶ 开始构建", width=120, height=36,
                                        fg_color="#007191", command=self._build)
        self.build_btn.pack(side="left", padx=20)
        self.open_btn = ctk.CTkButton(btn_frame, text="打开文件夹", width=100, command=self._open_folder)
        self.open_btn.pack(side="right", padx=5)

        # ── 输出位置选择 ──
        out_frame = ctk.CTkFrame(self)
        out_frame.pack(fill="x", padx=10, pady=5)
        ctk.CTkLabel(out_frame, text="输出位置:", width=70).pack(side="left", padx=5)
        self.out_var = ctk.StringVar(value="")
        self.out_entry = ctk.CTkEntry(out_frame, textvariable=self.out_var, width=400)
        self.out_entry.pack(side="left", fill="x", expand=True, padx=5)
        ctk.CTkButton(out_frame, text="浏览", width=60, command=self._browse_output).pack(side="right", padx=5)

        # ── 日志输出区 ──
        ctk.CTkLabel(self, text="构建日志", font=ctk.CTkFont(size=14, weight="bold")).pack(anchor="w", padx=15, pady=(5, 0))
        self.log_text = ctk.CTkTextbox(self, height=350, font=ctk.CTkFont(family="Consolas", size=12))
        self.log_text.pack(fill="both", expand=True, padx=10, pady=5)

        # ── 状态 ──
        self.status = ctk.CTkLabel(self, text="就绪", text_color="gray")
        self.status.pack(anchor="w", padx=15, pady=5)

        # 设置默认输出路径
        self._set_default_output()

    def _set_default_output(self):
        """根据论文路径设置默认输出文件名。"""
        pdf_path = self.shared.get('pdf_path', '')
        if pdf_path:
            base = os.path.splitext(os.path.basename(pdf_path))[0]
            default = os.path.join(os.path.dirname(pdf_path), f'{base}_汇报.pptx')
        else:
            default = os.path.join(os.path.expanduser('~'), 'Desktop', 'output.pptx')
        self.out_var.set(default)
        self._last_dir = os.path.dirname(default)

    def _browse_output(self):
        """打开文件保存对话框选择输出位置。"""
        initial_dir = self._last_dir if hasattr(self, '_last_dir') else os.path.expanduser('~/Desktop')
        path = filedialog.asksaveasfilename(
            title="选择输出位置",
            initialdir=initial_dir,
            defaultextension=".pptx",
            filetypes=[("PowerPoint", "*.pptx"), ("All Files", "*.*")]
        )
        if path:
            self.out_var.set(path)
            self._last_dir = os.path.dirname(path)

    def _log(self, msg):
        """追加日志到文本区。"""
        self.log_text.insert("end", msg + "\n")
        self.log_text.see("end")
        self.update_idletasks()

    def _build(self):
        """开始构建流程：验证 → 构建。"""
        log_step('build', '开始构建')
        json_str = self.shared.get('slide_content_json', '')
        if not json_str:
            messagebox.showerror("错误", "请先在大纲页生成或编辑 slide-content.json")
            return

        try:
            config = json.loads(json_str)
        except json.JSONDecodeError as e:
            messagebox.showerror("JSON 格式错误", str(e))
            return

        # 获取用户选择的输出路径
        output_path = self.out_var.get().strip()
        if not output_path:
            messagebox.showerror("错误", "请指定输出位置")
            return

        # 安全检查：目录是否存在
        out_dir = os.path.dirname(output_path)
        if out_dir and not os.path.exists(out_dir):
            try:
                os.makedirs(out_dir, exist_ok=True)
                self._log(f"[目录] 已创建: {out_dir}")
            except OSError as e:
                messagebox.showerror("错误", f"无法创建输出目录:\n{e}")
                return

        # 安全检查：文件是否已存在
        if os.path.exists(output_path):
            ok = messagebox.askyesno("文件已存在", f"{os.path.basename(output_path)} 已存在。\n是否覆盖？")
            if not ok:
                return

        # 检查目录是否可写
        if out_dir and not os.access(out_dir, os.W_OK):
            messagebox.showerror("错误", f"输出目录不可写:\n{out_dir}")
            return

        # 设置路径到配置
        config.setdefault('meta', {})
        config['meta']['template_path'] = self.shared.get('template_path', '')
        config['meta']['figs_dir'] = self.shared.get('figs_dir', '')
        config['meta']['output_path'] = output_path
        config['meta']['paper_meta'] = self.shared.get('paper_meta', {})

        # 写入临时 JSON 文件
        tmp = tempfile.NamedTemporaryFile(mode='w', suffix='.json', delete=False, encoding='utf-8')
        json.dump(config, tmp, ensure_ascii=False, indent=2)
        tmp_path = tmp.name
        tmp.close()

        self.output_path = output_path
        log_step('build', f'开始构建 → {output_path}')
        self._log(f"[开始] 输出 → {output_path}")
        self.build_btn.configure(state="disabled")
        self.status.configure(text="构建中...", text_color="gray")
        threading.Thread(target=self._do_build, args=(tmp_path,), daemon=True).start()

    def _do_build(self, json_path):
        """在后台线程中执行构建。"""
        import time
        t0 = time.time()
        try:
            self._log("[验证] 检查并修复 JSON...")
            sys.path.insert(0, os.path.join(os.path.dirname(__file__), '..', '..', 'scripts'))
            from validate_outline import repair_and_validate
            with open(json_path, 'r', encoding='utf-8') as f:
                raw = f.read()
            repaired_str, warnings = repair_and_validate(raw)
            if repaired_str is None:
                raise RuntimeError(f"JSON 无法修复: {warnings[0] if warnings else '未知错误'}")
            for w in warnings:
                self._log(f"[修复] {w}")
            if warnings:
                with open(json_path, 'w', encoding='utf-8') as f:
                    f.write(repaired_str)

            self._log("[构建] 开始生成 PPT...")
            from ppt_builder import load_json, build
            config = load_json(json_path)
            output = build(config, json_path)
            self.output_path = output
            elapsed = time.time() - t0
            log_step('build', f'构建完成 ({elapsed:.1f}s) → {output}')
            self._log(f"[完成] ✓ {output}")
            # 边界检测
            overflow = self._check_overflow(output)
            if overflow:
                self._log(f"[警告] {len(overflow)} 处元素超出边界:")
                for msg in overflow[:10]:
                    self._log(f"  {msg}")
            self.status.configure(text=f"✓ 构建完成 → {output}", text_color="green")
        except Exception as e:
            log_error(f"构建失败: {e}", exc_info=True)
            self._log(f"[错误] {e}")
            self.status.configure(text=f"✗ 构建失败: {e}", text_color="red")
        finally:
            self.build_btn.configure(state="normal")
            try:
                os.unlink(json_path)
            except:
                pass

    def _check_overflow(self, pptx_path):
        """检测内容页中文本框是否超出其背景框范围。返回警告列表。"""
        warnings = []
        try:
            from pptx import Presentation as Prs
            prs = Prs(pptx_path)
            from scripts.make_template import classify
            for si, slide in enumerate(prs.slides):
                cat = classify(slide)
                if cat not in ('CONTENT', 'CONTENT_FRAME', 'DISCUSSION'):
                    continue
                # 收集背景框(浅色大矩形)和文本框
                bg_boxes = []  # (name, x, y, w, h)
                text_boxes = []
                for sh in slide.shapes:
                    x = sh.left; y = sh.top; w = sh.width; h = sh.height
                    is_bg = False
                    try:
                        if sh.fill.type is not None:
                            rgb = sh.fill.fore_color.rgb
                            if min(rgb[0], rgb[1], rgb[2]) > 200:  # 浅色=背景框
                                is_bg = True
                    except: pass
                    if is_bg and w > 914400 * 0.5:  # >0.5in宽
                        bg_boxes.append((sh.name, x, y, x+w, y+h))
                    if sh.has_text_frame and sh.text_frame.text.strip():
                        text_boxes.append((sh.name, x, y, x+w, y+h, sh.text_frame.text.strip()[:30]))
                # 检查每个文本框是否在其下方的背景框内
                for t_name, tx1, ty1, tx2, ty2, txt in text_boxes:
                    inside = False
                    for b_name, bx1, by1, bx2, by2 in bg_boxes:
                        if bx1 - 5000 <= tx1 and by1 - 5000 <= ty1 and bx2 + 5000 >= tx2 and by2 + 5000 >= ty2:
                            inside = True; break
                    if not inside and ty1 < 6858000 * 0.9:  # 忽略页脚区域
                        y_in = ty1 / 914400; h_in = (ty2 - ty1) / 914400
                        warnings.append(f"Slide{si} {t_name}: 文字\"{txt}\" @ y={y_in:.1f}in h={h_in:.2f}in 无背景框")
        except Exception as e:
            warnings.append(f"边界检测失败: {e}")
        return warnings

    def _open_folder(self):
        """打开生成的 PPTX 文件（优先）或其所在文件夹。"""
        path = self.output_path or self.out_var.get().strip()
        if not path:
            messagebox.showinfo("提示", "尚未构建")
            return

        # 安全检查：确保路径在合理范围内
        norm = os.path.normpath(os.path.abspath(path))
        if not norm.lower().endswith('.pptx'):
            messagebox.showinfo("提示", "输出文件必须是 .pptx 格式")
            return

        if os.path.isfile(norm):
            # 文件已生成 → 直接打开
            try:
                os.startfile(norm)
                self._log(f"[打开] {norm}")
            except OSError as e:
                messagebox.showerror("错误", f"无法打开文件:\n{e}")
        else:
            # 文件尚未生成 → 打开所在文件夹
            out_dir = os.path.dirname(norm)
            if os.path.isdir(out_dir):
                os.startfile(out_dir)
            else:
                messagebox.showinfo("提示", "输出目录不存在，请先构建")
