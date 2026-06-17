#!/usr/bin/env python3
"""
ppt_master_adapter.py — ppt-master 适配层。

封装 ppt-master 的 template_fill_pptx 和 pptx_to_svg 管道，
为 PaperDeck 提供非破坏性的模板分析和设计 token 提取。

用法:
    from scripts.ppt_master_adapter import (
        is_available,           # ppt-master 是否可用
        analyze_template,       # 综合分析模板 → {library, theme_colors, theme_fonts, ...}
        extract_design_tokens,  # 提取增强设计 token → template_design.json 兼容格式
        detect_is_template,     # 判断是否已是模板骨架
    )

回退: 当 ppt-master 不可用时，is_available() 返回 False，调用方应回退到规则模式。
"""

import sys
import os
import json
from pathlib import Path

# ═══════════════════════════════════════════
# ppt-master 可用性检测 & 导入
# ═══════════════════════════════════════════

_PPTMASTER_SCRIPTS = None
_PPTMASTER_AVAILABLE = False
_LAST_ERROR = None  # 供调试


def _find_pptmaster_dir():
    """
    查找 ppt-master 核心模块目录。
    优先级: 内置 vendor/ > 外部 clone
    """
    candidates = []

    # 1. 内置 vendored 模块（始终可用）
    try:
        me = Path(__file__).resolve()
        vendor_dir = me.parent / 'vendor' / 'pptmaster'
        candidates.append(vendor_dir)
    except Exception:
        pass

    # 2. 外部 ppt-master clone: PaperDeck/../ppt-master/skills/ppt-master/scripts/
    try:
        me = Path(__file__).resolve()
        paperdeck_root = me.parent.parent  # PaperDeck/
        candidates.append(paperdeck_root.parent / 'ppt-master' / 'skills' / 'ppt-master' / 'scripts')
    except Exception:
        pass

    # 3. CWD 推导
    for base in [Path.cwd(), Path.cwd().parent]:
        candidates.append(base / 'ppt-master' / 'skills' / 'ppt-master' / 'scripts')

    # 4. 环境变量
    for env in ['PPTMASTER_HOME']:
        val = os.environ.get(env, '')
        if val:
            candidates.append(Path(val) / 'skills' / 'ppt-master' / 'scripts')

    # 返回第一个存在的（且包含 template_fill_pptx）
    for scripts_dir in candidates:
        if scripts_dir.exists() and (scripts_dir / 'template_fill_pptx').exists():
            return scripts_dir

    return None


def _init_pptmaster():
    """延迟初始化 ppt-master 导入路径。"""
    global _PPTMASTER_SCRIPTS, _PPTMASTER_AVAILABLE, _LAST_ERROR

    if _PPTMASTER_AVAILABLE:
        return True

    scripts_dir = _find_pptmaster_dir()
    if scripts_dir is None:
        _LAST_ERROR = "未找到 ppt-master 模块（内置 vendor/ 或外部 clone）"
        return False

    _PPTMASTER_SCRIPTS = scripts_dir
    if str(scripts_dir) not in sys.path:
        sys.path.insert(0, str(scripts_dir))

    try:
        import types
        # 注册 template_fill_pptx 包
        tfp_dir = scripts_dir / 'template_fill_pptx'
        if 'template_fill_pptx' not in sys.modules:
            pkg = types.ModuleType('template_fill_pptx')
            pkg.__path__ = [str(tfp_dir)]
            sys.modules['template_fill_pptx'] = pkg

        # 注册 pptx_to_svg 包
        pts_dir = scripts_dir / 'pptx_to_svg'
        if 'pptx_to_svg' not in sys.modules:
            pkg = types.ModuleType('pptx_to_svg')
            pkg.__path__ = [str(pts_dir)]
            sys.modules['pptx_to_svg'] = pkg

        _PPTMASTER_AVAILABLE = True
        _LAST_ERROR = None
        return True
    except Exception as e:
        _LAST_ERROR = str(e)
        return False


def is_available() -> bool:
    """ppt-master 是否可导入。"""
    return _init_pptmaster()


# ═══════════════════════════════════════════
# 公共 API
# ═══════════════════════════════════════════


def analyze_template(pptx_path: str) -> dict:
    """
    综合分析模板，返回统一的分析结果。

    内部调用:
      - analyze_pptx()  → slide_library (所有 text slot/table/chart)
      - convert_pptx_to_svg() → theme_colors + theme_fonts

    参数:
        pptx_path: 模板 PPTX 文件路径

    返回:
        {
            "library": {...},            # slide_library 完整字典
            "theme_colors": {...},       # {dk1, lt1, accent1-6, hlink, folHlink}
            "theme_fonts": {...},        # {majorLatin, minorLatin, ...}
            "slide_count": int,
            "canvas_px": (w, h),
            "page_types": {slide_index: page_type},  # {1: "cover_candidate", ...}
        }

    异常: 如果 ppt-master 不可用或分析失败，抛出 RuntimeError。
    """
    if not _init_pptmaster():
        raise RuntimeError("ppt-master 不可用，请确认已 clone 到 PaperDeck 同级目录")

    from template_fill_pptx.analyzer import analyze_pptx
    from pptx_to_svg.converter import convert_pptx_to_svg, ConvertOptions

    src = Path(pptx_path).resolve()
    if not src.exists():
        raise FileNotFoundError(f"模板文件不存在: {pptx_path}")

    # Step 1: 结构分析 (slide_library)
    try:
        library = analyze_pptx(src)
    except Exception as e:
        raise RuntimeError(f"模板结构分析失败: {e}") from e

    # Step 2: 主题提取 (theme colors + fonts)
    theme_colors = {}
    theme_fonts = {}
    canvas_px = (1280, 720)
    try:
        opts = ConvertOptions(inheritance_mode='flat')
        svg_result = convert_pptx_to_svg(src, options=opts)
        theme_colors = svg_result.theme_colors
        theme_fonts = svg_result.theme_fonts
        canvas_px = svg_result.canvas_px
    except Exception as e:
        # 主题提取失败不致命，继续使用 slide_library
        print(f"  [ppt-master] 主题提取失败 ({e})，仅使用结构分析")

    # Step 3: 汇总 page_types
    page_types = {}
    for slide in library.get('slides', []):
        idx = slide.get('slide_index')
        ptype = slide.get('page_type', 'content_candidate')
        if idx is not None:
            page_types[idx] = ptype

    return {
        'library': library,
        'theme_colors': theme_colors,
        'theme_fonts': theme_fonts,
        'slide_count': library.get('slide_count', 0),
        'canvas_px': canvas_px,
        'page_types': page_types,
    }


def extract_design_tokens(template_path: str, analysis: dict = None) -> dict:
    """
    从模板提取增强的设计 token，输出兼容 template_design.json 的格式。

    优先使用 ppt-master 的主题数据；回退到现有 template_extractor.extract()。

    参数:
        template_path: 模板 PPTX 路径
        analysis:      可选的预分析结果（避免重复分析）

    返回:
        dict, 兼容 template_design.json，额外增加 "theme" 段和 "_source" 标记
    """
    # 先尝试用 ppt-master
    if _init_pptmaster():
        try:
            if analysis is None:
                analysis = analyze_template(template_path)

            # ppt_master_design.py 和本文件同目录，直接 import
            try:
                from scripts.ppt_master_design import build_design_json_from_theme
            except ImportError:
                # 命令行直接运行时 scripts/ 不是 package，用同级 import
                import ppt_master_design
                build_design_json_from_theme = ppt_master_design.build_design_json_from_theme

            design = build_design_json_from_theme(
                analysis.get('theme_colors', {}),
                analysis.get('theme_fonts', {}),
                analysis.get('library', {}),
            )

            # 补充：ppt-master 无法提取非文本装饰形状（页眉/页脚条），
            # 用 template_extractor 补齐。需确保 scripts/ 在 sys.path 中
            # （template_extractor 内部会 from make_template import classify）
            import os as _os
            _script_dir = str(Path(__file__).resolve().parent)
            _already_in_path = _script_dir in sys.path
            if not _already_in_path:
                sys.path.insert(0, _script_dir)
            try:
                import template_extractor
                legacy = template_extractor.extract(template_path)
                # 装饰形状（非文本，ppt-master 不提取）
                if legacy.get('header_decorations'):
                    design['header_decorations'] = legacy['header_decorations']
                if legacy.get('footer_bars'):
                    design['footer_bars'] = legacy['footer_bars']
                # 字体/颜色：rule 提取器从实际 text run 读取，比 theme 推导更准
                for key in ('content_title', 'content_body', 'cover_title',
                            'section_number', 'section_title', 'toc_number', 'toc_item'):
                    if legacy.get(key) and legacy[key].get('font_name'):
                        design[key]['font_name'] = legacy[key]['font_name']
                    if legacy.get(key) and legacy[key].get('color'):
                        design[key]['color'] = legacy[key]['color']
                    if legacy.get(key) and legacy[key].get('font_size_pt'):
                        design[key]['font_size_pt'] = legacy[key]['font_size_pt']
            except Exception:
                pass  # 装饰提取失败不致命
            finally:
                if not _already_in_path:
                    sys.path.remove(_script_dir)

            return design
        except Exception as e:
            print(f"  [ppt-master] 设计 token 提取失败 ({e})，回退到规则方式")

    # 回退到现有 template_extractor
    try:
        from scripts.template_extractor import extract as extract_design
    except ImportError:
        import template_extractor
        extract_design = template_extractor.extract

    try:
        design = extract_design(template_path)
        design['_source'] = 'rule'
        return design
    except Exception as e:
        print(f"  [template_extractor] 提取失败 ({e})，返回空设计")
        return {'_source': 'fallback', '_error': str(e)}


def detect_is_template(pptx_path: str) -> bool:
    """
    判断 PPTX 是否已经是模板骨架（非完整内容 PPTX）。

    基于 slide_library 分析：
      - content 类型 slide 的文本过半包含占位符模式 ("此处填充", "Paper Title", "CONTENTS", 纯数字等)
      - 或者 slide 数量 <= 5（说明已是去重后的骨架）

    参数:
        pptx_path: 模板 PPTX 文件路径

    返回:
        True 表示已是骨架模板，False 表示是完整内容 PPTX
    """
    if not _init_pptmaster():
        # 回退：用规则 classify
        try:
            from pptx import Presentation
            from scripts.make_template import classify
            prs = Presentation(pptx_path)
            c_count = 0
            placeholder_count = 0
            for slide in prs.slides:
                cat = classify(slide)
                if cat in ('CONTENT', 'CONTENT_FRAME', 'DATA', 'DISCUSSION'):
                    c_count += 1
                    # 检查占位符文本
                    for shape in slide.shapes:
                        if shape.has_text_frame:
                            txt = shape.text_frame.text
                            if any(kw in txt for kw in ('此处填充', '章节标题', '论文标题')):
                                placeholder_count += 1
                                break
            if c_count > 0:
                return placeholder_count / c_count > 0.5
            return True  # 没有内容页 → 就是骨架
        except Exception:
            return False

    try:
        analysis = analyze_template(pptx_path)
        library = analysis.get('library', {})
        slides = library.get('slides', [])

        # 如果 slide <= 5，大概率是骨架
        if len(slides) <= 5:
            return True

        # 检查 content 类型 slide 的文本是否包含占位符模式
        placeholder_keywords = (
            '此处填充', '章节标题', '论文标题', 'Paper Title',
            '汇报人', 'CONTENTS', 'THANKS',
        )
        content_slides = [s for s in slides
                          if s.get('page_type') in ('content_candidate', 'chapter_candidate')]
        if not content_slides:
            return True

        placeholder_hits = 0
        for slide in content_slides:
            for slot in slide.get('slots', []):
                text = slot.get('text', '')
                if any(kw.lower() in text.lower() for kw in placeholder_keywords):
                    placeholder_hits += 1
                    break
                # 纯数字（01, 02...）也算占位符
                if text.strip().isdigit() and len(text.strip()) <= 2:
                    placeholder_hits += 1
                    break

        ratio = placeholder_hits / len(content_slides)
        return ratio > 0.5
    except Exception:
        return False


# ═══════════════════════════════════════════
# 便利函数: 保存/加载分析缓存
# ═══════════════════════════════════════════


def save_analysis_cache(analysis: dict, cache_path: str) -> None:
    """将分析结果保存为 JSON 缓存文件。"""
    # 移除不可序列化的对象
    serializable = {
        'library': analysis.get('library', {}),
        'theme_colors': analysis.get('theme_colors', {}),
        'theme_fonts': analysis.get('theme_fonts', {}),
        'slide_count': analysis.get('slide_count', 0),
        'page_types': analysis.get('page_types', {}),
    }
    with open(cache_path, 'w', encoding='utf-8') as f:
        json.dump(serializable, f, indent=2, ensure_ascii=False)


def load_analysis_cache(cache_path: str) -> dict | None:
    """从缓存加载分析结果。如果缓存不存在或过旧，返回 None。"""
    if not os.path.exists(cache_path):
        return None
    try:
        with open(cache_path, 'r', encoding='utf-8') as f:
            return json.load(f)
    except (json.JSONDecodeError, IOError):
        return None


# ═══════════════════════════════════════════
# CLI 测试入口
# ═══════════════════════════════════════════

if __name__ == '__main__':
    import sys as _sys
    if len(_sys.argv) < 2:
        print("用法: python ppt_master_adapter.py <template.pptx>")
        _sys.exit(1)

    path = _sys.argv[1]
    print(f"=== 分析模板: {path} ===")
    print(f"ppt-master 可用: {is_available()}")

    try:
        result = analyze_template(path)
        print(f"\n幻灯片数: {result['slide_count']}")
        print(f"页面类型分布:")
        for idx in sorted(result['page_types']):
            print(f"  第{idx}页: {result['page_types'][idx]}")

        print(f"\n主题色 ({len(result['theme_colors'])} 项):")
        for k, v in result['theme_colors'].items():
            print(f"  {k}: {v}")

        print(f"\n主题字体:")
        for k, v in result['theme_fonts'].items():
            print(f"  {k}: {v}")

        print(f"\n=== 设计 token 提取 ===")
        design = extract_design_tokens(path, result)
        print(f"来源: {design.get('_source', 'unknown')}")
        for key in design:
            if key.startswith('_'):
                continue
            val = design[key]
            if isinstance(val, dict) and 'font_name' in val:
                print(f"  {key}: font={val.get('font_name')}, size={val.get('font_size_pt')}pt, color=#{val.get('color', '?')}")
            elif isinstance(val, list):
                print(f"  {key}: [{len(val)} 项]")
            else:
                print(f"  {key}: {str(val)[:60]}")

        print(f"\n是否已是模板: {detect_is_template(path)}")

    except Exception as e:
        print(f"错误: {e}")
        import traceback
        traceback.print_exc()
        _sys.exit(1)
