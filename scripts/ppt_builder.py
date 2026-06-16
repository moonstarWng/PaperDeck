#!/usr/bin/env python3
"""
ppt_builder.py — paper2ppt 的 JSON 驱动 PPT 组装器。
读取 slide-content.json，加载模板 PPTX，构建所有幻灯片，
编辑模板页（封面/目录/章节/致谢），按计划重排序，保存最终 PPTX。

用法: python ppt_builder.py <slide-content.json>
"""
import copy, json, os, sys
from lxml import etree
from pptx.enum.text import PP_ALIGN
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# 将 scripts 目录加入路径，支持同目录导入
sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))
from ppt_layout import (
    FONT_EN, FONT_CN, BODY_SIZE, TEAL, DARK, WHITE,
    blank_layout, make_result_slide, parse_color, init_design,
)
from ppt_slides import (
    build_author_slide, build_background_slide, build_summary_slide,
    build_discussion1_slide, build_discussion2_slide, build_paper_info_slide,
)

# PPTX 内部 XML 命名空间，用于直接操作幻灯片列表和文本元素
NS_P = 'http://schemas.openxmlformats.org/presentationml/2006/main'
NS_A = 'http://schemas.openxmlformats.org/drawingml/2006/main'


def load_json(path):
    """从文件加载 JSON 配置。"""
    with open(path, 'r', encoding='utf-8') as f:
        return json.load(f)


def resolve_path(json_path, relative_path):
    """
    将 slide-content.json 中的相对路径解析为绝对路径。
    所有路径相对于 JSON 文件所在目录解析，避免工作目录依赖。
    """
    if os.path.isabs(relative_path):
        return relative_path
    json_dir = os.path.dirname(os.path.abspath(json_path))
    return os.path.normpath(os.path.join(json_dir, relative_path))


# ═══════════════════════════════════════════
# 模板幻灯片编辑器（仅修改文字，不修改任何图形元素）
# ═══════════════════════════════════════════

def edit_cover(slide, cover_data, tmpl_path=''):
    """
    修改封面幻灯片文字。
    通过位置特征识别标题区域和汇报人区域，替换文字内容。

    封面模板特征：
      - 标题区：垂直位置 2.0-3.0in，高度 > 1.0in（大号文字框）
      - 汇报人区：垂直位置 5.0-6.0in，水平位置 > 5.0in，宽度 < 4.0in（右下角小文字框）

    注意：汇报人区的原始宽度仅约 2.1in，需扩大至 4.5in 才能容纳完整信息。
    """
    title_en = cover_data.get('title_en', '')
    presenter = cover_data.get('presenter', 'xxx')
    date = cover_data.get('date', '')

    # 收集所有文本形状（用于智能识别标题和汇报人位置）
    text_shapes = []
    for shape in slide.shapes:
        if shape.has_text_frame and shape.text_frame.text.strip():
            t = shape.top / 914400
            l = shape.left / 914400
            w = shape.width / 914400
            h = shape.height / 914400
            text_shapes.append((shape, t, l, w, h, shape.text_frame.text.strip()))

    if not text_shapes:
        return

    # 标题 = 面积最大的文本框（宽度 > 3in 且高度 > 0.5in）
    candidates = [s for s in text_shapes if s[3] > 3 and s[4] > 0.5]
    title_shape = max(candidates, key=lambda s: s[3] * s[4]) if candidates else max(text_shapes, key=lambda s: len(s[5]))

    # 汇报人 = 含关键词 或 下半部分小文本框
    presenter_shape = None
    for s in text_shapes:
        txt = s[5].lower()
        if '汇报人' in txt or 'presenter' in txt:
            presenter_shape = s
            break
    if not presenter_shape:
        small = [s for s in text_shapes if s[1] > 3.5 and s[3] < 6 and s[4] < 1.0 and s != title_shape]
        if small:
            presenter_shape = min(small, key=lambda s: s[3] * s[4])

    # 标题字体：尝试从 design.json 读取 cover_title
    title_font = FONT_EN
    title_sz = Pt(32)
    try:
        import json as _json
        dj = tmpl_path.replace('.pptx', '_design.json')
        if os.path.exists(dj):
            with open(dj, 'r', encoding='utf-8') as f:
                ct = _json.load(f).get('cover_title', {})
                if ct.get('font_name'): title_font = ct['font_name']
                if ct.get('font_size_pt'): title_sz = Pt(ct['font_size_pt'])
    except: pass

    # 替换标题文字
    tf = title_shape[0].text_frame; tf.clear()
    for i, line in enumerate(title_en.split('\n')):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        r = p.add_run(); r.text = line
        r.font.name = title_font; r.font.size = title_sz; r.font.bold = True

    # 替换汇报人文字
    if presenter_shape:
        ps = presenter_shape[0]
        if ps.width / 914400 < 4.0:
            ps.width = Inches(4.5)
        tf = ps.text_frame; tf.clear()
        p = tf.paragraphs[0]
        r1 = p.add_run(); r1.text = '汇报人：'
        r1.font.name = FONT_EN; r1.font.size = Pt(20); r1.font.bold = True
        r2 = p.add_run(); r2.text = f'{presenter}    '
        r2.font.name = FONT_CN; r2.font.size = Pt(20); r2.font.bold = True
        if date:
            r3 = p.add_run(); r3.text = date
            r3.font.name = FONT_EN; r3.font.size = Pt(16)


def edit_toc(slide, toc_map, section_titles=None):
    """
    修改目录页文字。
    toc_map: {旧文本: 新文本} 替换映射。为空时若提供 section_titles 则自动按顺序替换。
    """

    # 自动构建映射：无 toc_map 但有 section_titles 时，按顺序替换 TOC 文本
    if not toc_map and section_titles:
        # 收集 TOC 条目对 (number_shape, title_shape) — 按位置识别
        num_shapes = []   # X < 2.0in → 编号
        title_shapes = [] # X >= 2.0in → 标题
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            if 'TEXT_BOX' not in str(shape.shape_type):
                continue
            txt = shape.text_frame.text.strip()
            x = shape.left / 914400
            if txt.upper() in ('CONTENTS', '目录', '目 录'):
                continue
            if x < 2.0:
                num_shapes.append(shape)
            else:
                title_shapes.append(shape)
        num_shapes.sort(key=lambda s: s.top)
        title_shapes.sort(key=lambda s: s.top)

        n_existing = min(len(num_shapes), len(title_shapes))
        n_needed = len(section_titles)

        # 收集分隔线
        import copy
        sep_lines = []
        for sh in slide.shapes:
            if 'AUTO_SHAPE' in str(sh.shape_type):
                x = sh.left / 914400; w = sh.width / 914400
                if x > 2.0 and w > 5 and sh.height / 914400 < 0.1:
                    sep_lines.append(sh)
        sep_lines.sort(key=lambda s: s.top)
        first_title_bottom = title_shapes[0].top + title_shapes[0].height
        first_sep = None
        for sl in sep_lines:
            if sl.top >= first_title_bottom - Inches(0.05):
                first_sep = sl; break

        title_dy = int(title_shapes[0].top - num_shapes[0].top)
        sep_dy = int(first_sep.top - num_shapes[0].top) if first_sep else int(Inches(0.55))
        total_rows = n_needed

        # ── 补充不足的条目 ──
        for pi in range(n_existing, total_rows):
            ns = _clone_shape(slide, num_shapes[-1]); num_shapes.append(ns)
            ts = _clone_shape(slide, title_shapes[-1]); title_shapes.append(ts)
            if first_sep:
                _clone_shape(slide, first_sep)

        # ── 重新收集所有行和分隔线 ──
        all_num = list(num_shapes)
        all_title = list(title_shapes)
        sep_lines2 = []
        for sh in slide.shapes:
            if 'AUTO_SHAPE' in str(sh.shape_type):
                x = sh.left / 914400; w = sh.width / 914400
                if x > 2.0 and w > 5 and sh.height / 914400 < 0.1:
                    sep_lines2.append(sh)
        sep_lines2.sort(key=lambda s: s.top)

        # ── 删除多余的条目（模板有 4 行但只需要 2 行时）──
        for pi in range(total_rows, len(all_num)):
            try:
                all_num[pi]._element.getparent().remove(all_num[pi]._element)
            except Exception: pass
        for pi in range(total_rows, len(all_title)):
            try:
                all_title[pi]._element.getparent().remove(all_title[pi]._element)
            except Exception: pass
        all_num = all_num[:total_rows]
        all_title = all_title[:total_rows]

        # ── 分隔线：补到 total_rows，多余的删 ──
        while len(sep_lines2) < total_rows:
            if first_sep:
                sl = _clone_shape(slide, first_sep)
                sep_lines2.append(sl)
            else:
                break
        for pi in range(total_rows, len(sep_lines2)):
            try:
                sep_lines2[pi]._element.getparent().remove(sep_lines2[pi]._element)
            except Exception: pass
        sep_lines2 = sep_lines2[:total_rows]

        # ── 固定高度均匀分布（总高度 = 原模板 4 行跨度 3.30in）──
        y_top = int(num_shapes[0].top)
        y_span = int(Inches(3.30))  # 固定总高，不随行数变
        y_last = y_top + y_span
        if total_rows > 1:
            row_gap = int((y_last - y_top) / max(total_rows - 1, 1))
        else:
            row_gap = 0

        for pi in range(total_rows):
            base_y = y_top + pi * row_gap
            all_num[pi].top = base_y
            all_title[pi].top = base_y + title_dy
        for pi in range(total_rows):
            if pi < len(sep_lines2):
                sep_lines2[pi].top = all_num[pi].top + sep_dy

        # ── 更新文字 ──
        for pi in range(total_rows):
            _set_shape_text(all_num[pi], f'0{pi+1}')
            _set_shape_text(all_title[pi], section_titles[pi])
        return

    if not toc_map:
        return
    for t_elem in slide._element.iter(f'{{{NS_A}}}t'):
        if t_elem.text:
            for old, new in toc_map.items():
                if old in t_elem.text:
                    t_elem.text = t_elem.text.replace(old, new)


def _set_shape_text(shape, text):
    """替换 shape 中所有文本为指定文字。"""
    if shape.has_text_frame:
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                run.text = ''
        shape.text_frame.paragraphs[0].runs[0].text = text if shape.text_frame.paragraphs[0].runs else ''
        if not shape.text_frame.paragraphs[0].runs:
            shape.text_frame.paragraphs[0].add_run().text = text


def _clone_shape(slide, src_shape):
    """在 slide 上克隆一个 shape（XML 深拷贝），返回新的 shape 对象。"""
    import copy
    new_el = copy.deepcopy(src_shape._element)
    slide.shapes._spTree.append(new_el)
    # 返回新 shape（通过名称匹配）
    for s in slide.shapes:
        if s._element is new_el:
            return s
    return None


def edit_section_divider(slide, number, title):
    """
    修改章节分隔页的编号和标题文字。
    编号特征：恰好2位数字
    标题特征：包含中文字符
    """
    for shape in slide.shapes:
        if shape.has_text_frame:
            text = shape.text_frame.text.strip()
            # 匹配两位数字编号（如 '01', '02'）
            if text.isdigit() and len(text) == 2:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if run.text.strip().isdigit() and len(run.text.strip()) == 2:
                            run.text = number
            # 匹配章节标题：中文或英文（只要不是纯数字且长度≥2）
            elif len(text) >= 2 and not text.isdigit():
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if not run.text.strip().isdigit() and len(run.text.strip()) >= 2:
                            run.text = title


# ═══════════════════════════════════════════
# 主编排器
# ═══════════════════════════════════════════

def build(config, json_path='.'):
    """
    主编排函数：加载模板 → 构建所有幻灯片 → 编辑模板页 → 重排序 → 保存。

    工作流程：
      1. 从 meta 中读取模板路径、输出路径、图片目录、模板幻灯片索引
      2. 按 slides 数组顺序构建所有内容幻灯片（追加到模板幻灯片末尾）
      3. 编辑模板幻灯片（封面/目录/章节分隔页）的文字
      4. 通过操作 sldIdLst XML 重排序所有幻灯片
      5. 保存最终 PPTX
    """
    meta = config.get('meta', {})
    # 将 JSON 中的相对路径解析为绝对路径
    template_path = resolve_path(json_path, meta.get('template_path', 'template.pptx'))
    output_path = resolve_path(json_path, meta.get('output_path', 'output.pptx'))
    figs_dir = resolve_path(json_path, meta.get('figs_dir', './figs'))
    indices = meta.get('template_slide_indices', {})

    # 加载模板设计令牌（缺失时当场生成到 process/）
    tmpl_dir = os.path.dirname(os.path.abspath(template_path))
    proc_dir = os.path.join(tmpl_dir, 'process')
    os.makedirs(proc_dir, exist_ok=True)
    design_json = os.path.join(proc_dir, 'template_design.json')
    if not os.path.exists(design_json):
        # 当场从模板提取
        try:
            from template_extractor import extract as extract_design
            design = extract_design(template_path)
            import json as _json
            os.makedirs(os.path.dirname(design_json), exist_ok=True)
            with open(design_json, 'w', encoding='utf-8') as f:
                _json.dump(design, f, indent=2, ensure_ascii=False)
            print(f"  自动生成设计令牌 → {design_json}")
        except Exception as e:
            fallback = template_path.replace('.pptx', '_design.json')
            if os.path.exists(fallback):
                design_json = fallback
            else:
                print(f"  WARNING: 设计令牌生成失败 ({e})，使用默认值")
    init_design(design_json)

    prs = Presentation(template_path)
    print(f"Template: {len(prs.slides)} slides")

    # 自动检测模板幻灯片索引（如果 JSON 中未指定或指定了不完整的）
    from make_template import classify
    auto_indices = {'cover': None, 'toc': 1, 'sections': [], 'thanks': len(prs.slides) - 1}
    for i, slide in enumerate(prs.slides):
        cat = classify(slide)
        if cat == 'COVER' and auto_indices['cover'] is None:
            auto_indices['cover'] = i   # 取第一个 COVER
        elif cat == 'TOC':
            auto_indices['toc'] = i
        elif cat.startswith('SECTION_'):
            auto_indices['sections'].append(i)
        elif cat == 'THANKS':
            auto_indices['thanks'] = i
    if auto_indices['cover'] is None:
        auto_indices['cover'] = 0  # fallback
    auto_indices['sections'].sort()
    # JSON 中的值覆盖自动检测值（允许用户手动指定）
    for key in ['cover', 'toc', 'thanks']:
        if key in indices:
            auto_indices[key] = indices[key]
    if 'sections' in indices and indices['sections']:
        auto_indices['sections'] = indices['sections']
    indices = auto_indices

    slides_plan = config.get('slides', [])
    section_divider_map = config.get('section_divider_edits', [])
    section_indices = list(indices.get('sections', []))

    # ── 章节页不够时用代码重建（必须在 new_order 构建之前）──
    def _add_section_slide(prs, base_slide):
        """重建一个与模板章节页布局一致的章节页。从模板复制字体样式。"""
        from pptx.dml.color import RGBColor
        # 从模板读取字体
        tmpl_num_font = None
        tmpl_title_font = None
        for sh in base_slide.shapes:
            if sh.has_text_frame:
                for p in sh.text_frame.paragraphs:
                    for r in p.runs:
                        if r.text.strip().isdigit() and len(r.text.strip()) == 2:
                            tmpl_num_font = r.font
                        elif len(r.text.strip()) >= 2:
                            tmpl_title_font = r.font

        layout = base_slide.slide_layout
        new_s = prs.slides.add_slide(layout)
        # 全幅深色背景（从模板读颜色）
        bg_color = RGBColor(0x1A, 0x1A, 0x1A)
        accent = RGBColor(0xE6, 0x39, 0x46)
        try:
            for sh in base_slide.shapes:
                if sh.width/914400 > 12 and sh.height/914400 > 6:
                    if sh.fill.type is not None:
                        bg_color = sh.fill.fore_color.rgb
                        break
        except: pass
        bg = new_s.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.33), Inches(7.5))
        bg.fill.solid(); bg.fill.fore_color.rgb = bg_color; bg.line.fill.background()
        # 左侧竖条
        bar = new_s.shapes.add_shape(1, Inches(0), Inches(0), Inches(1.2), Inches(7.5))
        bar.fill.solid(); bar.fill.fore_color.rgb = accent; bar.line.fill.background()
        # 编号 TextBox
        num_box = new_s.shapes.add_textbox(Inches(1.5), Inches(2.5), Inches(1.5), Inches(1.0))
        p = num_box.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
        r = p.add_run(); r.text = '00'
        if tmpl_num_font:
            r.font.size = tmpl_num_font.size; r.font.bold = tmpl_num_font.bold
            r.font.name = tmpl_num_font.name
            try: r.font.color.rgb = tmpl_num_font.color.rgb
            except: r.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        else:
            r.font.size = Pt(56); r.font.bold = True
            r.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF); r.font.name = 'Arial'
        # 标题 TextBox
        title_box = new_s.shapes.add_textbox(Inches(3.0), Inches(2.5), Inches(7.0), Inches(1.0))
        tp = title_box.text_frame.paragraphs[0]
        tr = tp.add_run(); tr.text = 'Section Title'
        if tmpl_title_font:
            tr.font.size = tmpl_title_font.size; tr.font.bold = tmpl_title_font.bold
            tr.font.name = tmpl_title_font.name
            try: tr.font.color.rgb = tmpl_title_font.color.rgb
            except: tr.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        else:
            tr.font.size = Pt(44); tr.font.bold = True
            tr.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF); tr.font.name = 'Arial'
        # 装饰线
        deco = new_s.shapes.add_shape(1, Inches(1.5), Inches(3.8), Inches(3.0), Pt(2))
        deco.fill.solid(); deco.fill.fore_color.rgb = accent; deco.line.fill.background()
        # 底部横条
        bot = new_s.shapes.add_shape(1, Inches(0), Inches(7.35), Inches(13.33), Pt(3))
        bot.fill.solid(); bot.fill.fore_color.rgb = accent; bot.line.fill.background()
        return new_s

    if section_indices and section_divider_map and len(section_indices) < len(section_divider_map):
        base_slide = prs.slides[section_indices[0]]
        n_need = len(section_divider_map) - len(section_indices)
        print(f"  模板仅 {len(section_indices)} 个章节页，自动创建 {n_need} 个")
        for _ in range(n_need):
            new_s = _add_section_slide(prs, base_slide)
            section_indices.append(len(prs.slides) - 1)

    new_order = []          # 最终幻灯片排列顺序（模板原始索引 + 新幻灯片索引）
    section_used = 0        # 跟踪下一次使用哪个章节分隔页

    # 幻灯片类型 → 构建函数的映射表
    slide_builders = {
        'author': build_author_slide,
        'background': build_background_slide,
        'summary': build_summary_slide,
        'discussion1': build_discussion1_slide,
        'discussion2': build_discussion2_slide,
        'paper_info': build_paper_info_slide,
    }

    # ── 遍历构建计划，逐一创建幻灯片 ──
    for item in slides_plan:
        stype = item.get('type')

        if stype == 'keep':
            # 保留模板原始幻灯片：按 ref 引用名映射到实际索引
            ref = item.get('ref')
            if ref == 'cover':
                new_order.append(indices.get('cover', 0))
            elif ref == 'toc':
                new_order.append(indices.get('toc', 1))
            elif ref == 'thanks':
                new_order.append(indices.get('thanks', 7))
            elif ref == 'section':
                idx = item.get('index', section_used)
                section_used = idx + 1
                if idx < len(section_indices):
                    new_order.append(section_indices[idx])
                else:
                    print(f"  WARNING: section index {idx} out of range")

        elif stype == 'result':
            # 结果内容页：标题 + 正文 + 图片
            title = item.get('title', '')
            body = item.get('body', [])
            imgs = item.get('images', [])
            make_result_slide(prs, title, body, imgs, figs_dir)
            new_order.append(len(prs.slides) - 1)
            print(f"  Result: {title[:60]}")

        elif stype in slide_builders:
            # 解析 paper_info 的相对路径
            if stype == 'paper_info' and item.get('pdf_path'):
                item['pdf_path'] = resolve_path(json_path, item['pdf_path'])
            slide_builders[stype](prs, item)
            new_order.append(len(prs.slides) - 1)
            print(f"  {stype}: {item.get('title', '')[:60]}")

        else:
            print(f"  WARNING: unknown slide type '{stype}'")

    print(f"\nBuilt {len(new_order)} slides in plan")

    # ── 编辑模板幻灯片文字 ──
    print("Editing template slides...")
    if 'cover' in config:
        edit_cover(prs.slides[indices.get('cover', 0)], config['cover'], template_path)
    if 'toc_replacements' in config:
        stitles = [sde['title'] for sde in section_divider_map] if section_divider_map else None
        edit_toc(prs.slides[indices.get('toc', 1)], config.get('toc_replacements', {}), stitles)
    # 编辑章节分隔页（已在前面克隆补充够数量）
    for i, sde in enumerate(section_divider_map):
        if i < len(section_indices):
            edit_section_divider(prs.slides[section_indices[i]], sde['number'], sde['title'])
    # 致谢页保持模板原文，不做任何修改

    # ── 重排序幻灯片 ──
    # python-pptx 不支持原生重排序，需直接操作 sldIdLst XML
    print("Reordering slides...")
    pres_elem = prs.part._element
    sldIdLst = pres_elem.find(f'{{{NS_P}}}sldIdLst')
    if sldIdLst is not None:
        all_ids = list(sldIdLst)
        # 清空当前顺序
        for c in list(sldIdLst):
            sldIdLst.remove(c)
        # 按 new_order 重新插入（使用 deepcopy 避免引用冲突）
        for i in new_order:
            if i < len(all_ids):
                sldIdLst.append(copy.deepcopy(all_ids[i]))
    print(f"Final: {len(prs.slides)} slides")

    # ── 保存 ──
    print(f"\nSaving {output_path}...")
    prs.save(output_path)
    print(f"DONE: {output_path}")
    return output_path


def main():
    """命令行入口：接收 slide-content.json 路径，执行完整构建流程。"""
    if len(sys.argv) < 2:
        print(f"Usage: python {sys.argv[0]} <slide-content.json>")
        sys.exit(1)
    config = load_json(sys.argv[1])
    build(config, sys.argv[1])


if __name__ == '__main__':
    main()
