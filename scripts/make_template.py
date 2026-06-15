#!/usr/bin/env python3
"""
make_template.py — 从完整的文献汇报 PPTX 中提取设计骨架模板。

处理步骤:
  1. 删除所有内容图片（仅保留 Logo 和章节图标）
  2. 删除包含嵌入图片的 GROUP 形状
  3. 替换汇报人姓名等隐私信息为 "xxx"
  4. 删除内容页（CONTENT / CONTENT_FRAME / DATA / DISCUSSION）
  5. 按页面类型去重（SECTION 全部保留，其余类型仅保留首个）

用法: python make_template.py <input.pptx> [output.pptx]
输出: 仅保留设计骨架的模板 PPTX（通常 7-10 页）
"""

import sys, os, re
from pptx import Presentation


# ═══════════════════════════════════════════
# Logo 和章节图标识别
# ═══════════════════════════════════════════

def _is_logo(shape):
    """
    判断一个 PICTURE 形状是否为 Logo。
    识别依据：位于右上角（x≈5.88in），尺寸约 1.51×1.44in。

    注意：这些阈值是根据特定模板校准的。换用不同模板时可能需要调整。
    """
    x, w, h = shape.left / 914400, shape.width / 914400, shape.height / 914400
    return abs(x - 5.88) < 0.3 and abs(w - 1.51) < 0.3 and abs(h - 1.44) < 0.3


def _is_icon(shape):
    """
    判断一个 PICTURE 形状是否为章节图标。
    识别依据：位于章节页深色横带上（x≈5.40in, y≈2.79in），尺寸约 0.78×0.78in。
    """
    x, y = shape.left / 914400, shape.top / 914400
    w, h = shape.width / 914400, shape.height / 914400
    return abs(x - 5.40) < 0.2 and abs(y - 2.79) < 0.2 and abs(w - 0.78) < 0.2 and abs(h - 0.78) < 0.2


# ═══════════════════════════════════════════
# Step 1: 清除内容图片和图片 GROUP
# ═══════════════════════════════════════════

def strip_images(prs):
    """
    遍历所有幻灯片，删除内容图片和包含图片的 GROUP 形状。
    保留 Logo 和章节图标。

    返回: (删除数, 保留数)
    """
    removed, kept = 0, 0
    for slide in prs.slides:
        # ── 删除 PICTURE 类型的形状（保留 Logo 和图标）──
        to_remove = []
        for shape in slide.shapes:
            if str(shape.shape_type) == 'PICTURE (13)':
                if _is_logo(shape) or _is_icon(shape):
                    kept += 1
                else:
                    to_remove.append(shape)
        for s in to_remove:
            s._element.getparent().remove(s._element)
            removed += 1

        # ── 删除包含 <p:pic> 子元素的 GROUP 形状 ──
        # GROUP 形状本身不是 PICTURE，但可能内嵌了 <p:pic> 元素
        # （例如：带图片的卡片组合）。用 lxml 的命名空间搜索子元素。
        ns = {'p': 'http://schemas.openxmlformats.org/presentationml/2006/main'}
        groups = []
        for shape in slide.shapes:
            if str(shape.shape_type) == 'GROUP (6)':
                if shape._element.findall('.//p:pic', ns):
                    groups.append(shape)
        for g in groups:
            g._element.getparent().remove(g._element)
            removed += 1
    return removed, kept


# ═══════════════════════════════════════════
# Step 2: 替换隐私信息
# ═══════════════════════════════════════════

def replace_name(prs, targets=None):
    """
    替换汇报人姓名等个人信息为 "xxx"。

    参数:
      targets: [(姓氏首字, 名字末字), ...] 列表。
               例如 [('张', '三')] 会匹配 "张...三" 中间可能含空格的姓名。
               默认为 None（不执行替换），使用时需要显式传入 targets。
    """
    if targets is None:
        # 不执行替换 —— 每个用户应传入自己的 targets
        return
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    for first, last in targets:
                        if first in run.text and last in run.text:
                            i1 = run.text.index(first)
                            i2 = run.text.index(last)
                            # 将 first...last 之间的全部字符替换为 xxx
                            run.text = run.text[:i1] + 'xxx' + run.text[i2 + 1:]


# ═══════════════════════════════════════════
# Step 3: 幻灯片分类
# ═══════════════════════════════════════════

# 内容页类型集合 —— 这些类型的页面将被全部删除
CONTENT_TYPES = {'CONTENT', 'CONTENT_FRAME', 'DATA', 'DISCUSSION'}


def classify(slide):
    """
    将幻灯片归类为以下类型之一:

      COVER        — 封面（FREEFORM 形状 + ≤3 个文本框）
      TOC          — 目录（≥2 个 GROUP，或含"目录"/"CONTENTS"）
      SECTION_XX   — 章节分隔页（FREEFORM + 章节编号 01-05）
      SUMMARY      — 小结页（含"\d+ 小结\d+"模式文本）
      THANKS       — 致谢页（含"THANKS"或"感谢"）
      DATA         — 数据页（≥3 个单字母标签 A/B/C...）
      DISCUSSION   — 讨论/总结页（含"总结"或"讨论"，不含"小结"）
      CONTENT_FRAME— 内容框架页（含圆角矩形框架）
      CONTENT      — 兜底内容页

    分类规则基于结构特征（形状类型、文本模式），不依赖硬编码页码，
    因此天然适配使用同一设计语言的不同 PPTX。
    """
    types = {}
    has_freeform = False
    has_dark_bg = False       # AUTO_SHAPE 矩形做深色背景（兼容非 FREEFORM 的模板）
    all_text = ''
    label_count = 0
    has_frame = False

    for shape in slide.shapes:
        st = str(shape.shape_type)
        types[st] = types.get(st, 0) + 1

        if st == 'FREEFORM (5)':
            has_freeform = True

        if st == 'AUTO_SHAPE (1)':
            name = shape.name.lower() if hasattr(shape, 'name') else ''
            if '圆角' in name or 'round' in name:
                has_frame = True
            # 检测全幅深色背景矩形（用于封面和章节页，替代 FREEFORM）
            w = shape.width / 914400
            h = shape.height / 914400
            if w > 10 or h > 5:
                has_dark_bg = True

        if shape.has_text_frame:
            t = shape.text_frame.text.strip()
            all_text += t + ' '
            if len(t) == 1 and t.isascii() and t.isalpha() and t == t.upper():
                label_count += 1

    # 只要是有大背景矩形或 FREEFORM，都视为"有装饰背景"
    has_background = has_freeform or has_dark_bg
    chapter_match = re.search(r'\b(\d{1,2})\b', all_text)

    # ── 按优先级分类（THANKS 和 TOC 优先于 SECTION/COVER）──
    if re.search(r'\d+\s*小结\d*', all_text):
        return 'SUMMARY'

    # THANKS 页必须最先检查（避免被误判为 COVER）
    if 'THANKS' in all_text.upper() or '感谢' in all_text:
        return 'THANKS'

    # TOC 优先于 SECTION（TOC 也可能包含章节编号 "01"）
    if types.get('GROUP (6)', 0) >= 2 or '目录' in all_text or 'CONTENTS' in all_text.upper():
        return 'TOC'

    if has_background and chapter_match and '目录' not in all_text:
        return f'SECTION_{chapter_match.group(1)}'

    if has_background and types.get('TEXT_BOX (17)', 0) <= 3:
        return 'COVER'

    if label_count >= 3:
        return 'DATA'

    if ('总结' in all_text or '讨论' in all_text) and '小结' not in all_text:
        return 'DISCUSSION'

    if has_frame:
        return 'CONTENT_FRAME'

    return 'CONTENT'


# ═══════════════════════════════════════════
# Step 4: 删除内容页 + 去重
# ═══════════════════════════════════════════

def trim_slides(prs):
    """
    删除所有 CONTENT / CONTENT_FRAME / DATA / DISCUSSION 页面。
    然后去重：保留所有 SECTION_XX 页面，每种其他类型仅保留首次出现的那个。

    通过操作 sldIdLst XML 删除幻灯片，同时删除对应的 relationship。
    返回: 删除的幻灯片数量。
    """
    to_keep = set()     # 要保留的幻灯片索引集合
    seen = {}           # 记录每种类型首次出现的索引

    for i, slide in enumerate(prs.slides):
        cat = classify(slide)

        # 内容页类型直接跳过
        if cat in CONTENT_TYPES:
            continue

        # 章节页全部保留，其余类型仅保留首个
        if cat.startswith('SECTION_'):
            to_keep.add(i)
        elif cat not in seen:
            seen[cat] = i
            to_keep.add(i)

    # 从后往前删除（避免索引偏移）
    to_delete = sorted([i for i in range(len(prs.slides)) if i not in to_keep], reverse=True)

    for i in to_delete:
        sldId = prs.slides._sldIdLst[i]
        rId = sldId.get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id')
        prs.slides._sldIdLst.remove(sldId)       # 从幻灯片列表中移除
        try:
            prs.part.drop_rel(rId)                 # 删除对应的关系引用
        except Exception:
            pass

    return len(to_delete)


# ═══════════════════════════════════════════
# Step 5: 通用化模板文字
# ═══════════════════════════════════════════

def generalize_text(prs):
    """
    将非内容页中的具体文字替换为通用占位符。
    封面标题 → "论文标题"，汇报人 → "汇报人：xxx"，
    章节标题 → "章节标题"，目录项 → 保持编号但清空描述。
    """
    for slide in prs.slides:
        cat = classify(slide)

        if cat == 'COVER':
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                text = shape.text_frame.text.strip()
                if not text:
                    continue

                # 多行文字、长文字 → 标题占位
                lines = text.split('\n')
                is_title = (len(lines) >= 2 or len(text) > 15)
                is_presenter = any(kw in text.lower() for kw in
                                   ('presenter', '汇报人', '报告人', '演讲人', 'xxx'))

                if is_presenter:
                    # 保留 "汇报人：" 前缀，人名改为 xxx
                    for para in shape.text_frame.paragraphs:
                        for run in para.runs:
                            if any(kw in run.text for kw in
                                   ('presenter', '汇报人', '报告人', '演讲人')):
                                run.text = '汇报人：'
                            elif len(run.text.strip()) > 0 and run.text.strip() != ':':
                                run.text = 'xxx'
                elif is_title:
                    # 大标题 → 占位
                    tf = shape.text_frame
                    tf.clear()
                    p = tf.paragraphs[0]
                    r = p.add_run()
                    r.text = '论文标题'
                    # 保留原字体大小
                    try:
                        r.font.size = shape.text_frame.paragraphs[0].runs[0].font.size
                    except:
                        pass

        elif cat.startswith('SECTION_'):
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                text = shape.text_frame.text.strip()
                if not text or text.isdigit() or (len(text) <= 2 and text.isascii()):
                    continue  # 保留数字编号
                # 非数字文本 → 章节标题占位
                tf = shape.text_frame
                tf.clear()
                p = tf.paragraphs[0]
                r = p.add_run()
                r.text = '章节标题'

        elif cat == 'TOC':
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        t = run.text.strip()
                        # 保留纯数字项 (01, 02...) 和 TOC 标题
                        if t.isdigit() and len(t) <= 2:
                            continue
                        if t.upper() in ('CONTENTS', '目录', '目 录'):
                            continue
                        # 其余替换
                        if len(t) > 2:
                            run.text = '章节名称'


# ═══════════════════════════════════════════
# 主入口
# ═══════════════════════════════════════════

def main():
    """命令行入口：读取完整 PPTX，输出模板骨架 PPTX。"""
    if len(sys.argv) < 2:
        print('Usage: python make_template.py <input.pptx> [output.pptx]')
        sys.exit(1)

    src = sys.argv[1]
    # 默认输出文件名：{输入文件名}_模板.pptx
    if len(sys.argv) >= 3:
        dst = sys.argv[2]
    else:
        base = os.path.splitext(os.path.basename(src))[0]
        dst = os.path.join(os.path.dirname(src) or '.', f'{base}_模板.pptx')

    if not os.path.exists(src):
        print(f'ERROR: {src} not found')
        sys.exit(1)

    prs = Presentation(src)
    n0 = len(prs.slides)

    r_img, k_img = strip_images(prs)
    print(f'  Content images removed: {r_img}, logos/icons kept: {k_img}')

    replace_name(prs)
    print(f'  Personal names replaced')

    d = trim_slides(prs)
    print(f'  Slides removed (content + duplicates): {d}')

    generalize_text(prs)
    print(f'  Template text generalized')

    prs.save(dst)
    n1 = len(prs.slides)
    kb = os.path.getsize(dst) / 1024
    print(f'  Slides: {n0} -> {n1}  |  Size: {kb:.0f} KB')
    print(f'  -> {dst}')


if __name__ == '__main__':
    main()
