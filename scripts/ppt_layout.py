"""
ppt_layout.py — paper2ppt 可复用布局工具库。
支持从 template_design.json 加载设计令牌，替代硬编码样式。
"""
import os, json
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ═══════════════════════════════════════════
# 默认颜色调色板（被设计令牌覆盖）
# ═══════════════════════════════════════════

TEAL = RGBColor(0x00, 0x71, 0x91)
DARK = RGBColor(0x33, 0x33, 0x33)
GRAY = RGBColor(0xA5, 0xA5, 0xA5)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
NAVY = RGBColor(0x1A, 0x2E, 0x4A)
LIGHT_TEAL = RGBColor(0xE0, 0xF0, 0xF5)
LIGHT_GREEN = RGBColor(0xE8, 0xF5, 0xE9)
RED_ACCENT = RGBColor(0xD9, 0x4F, 0x4F)
GREEN_ACCENT = RGBColor(0x2E, 0x7D, 0x32)
LIGHT_RED = RGBColor(0xFD, 0xED, 0xEC)
GOLD = RGBColor(0xD4, 0x8B, 0x2C)
ORANGE = RGBColor(0xE0, 0x7B, 0x3A)
BLUE_SLATE = RGBColor(0x2E, 0x86, 0xAB)
BLUE_DEEP = RGBColor(0x1B, 0x5E, 0x7A)

COLOR_MAP = {
    'teal': TEAL, 'dark': DARK, 'gray': GRAY, 'white': WHITE, 'navy': NAVY,
    'light_teal': LIGHT_TEAL, 'light_green': LIGHT_GREEN,
    'red_accent': RED_ACCENT, 'green_accent': GREEN_ACCENT, 'light_red': LIGHT_RED,
    'gold': GOLD, 'orange': ORANGE, 'blue_slate': BLUE_SLATE, 'blue_deep': BLUE_DEEP,
}

# ═══════════════════════════════════════════
# 设计令牌（默认值 → 被 init_design() 覆盖）
# ═══════════════════════════════════════════

FONT_EN = 'Times New Roman'
FONT_CN = '微软雅黑'
BODY_SIZE = Pt(16)
IMG_W = 5.5
ACCENT_COLOR = TEAL        # 主色：圆点、标题装饰
TITLE_COLOR = TEAL         # 内容页标题色
TITLE_SIZE = Pt(24)        # 内容页标题字号
BODY_COLOR = DARK          # 正文色
FOOTER_BARS = [(7.2, GRAY, 0.095), (7.3, TEAL, 0.2)]  # [(y, color, height)]
HEADER_DECOS = []          # 标题装饰 [(x, y, w, h, color)]


def _hex_to_rgb(h):
    """'#RRGGBB' 或 'RRGGBB' → RGBColor。"""
    h = h.lstrip('#')
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


# ═══════════════════════════════════════════
# 色板生成工具
# ═══════════════════════════════════════════

def _tint(rgb, factor=0.8):
    """与白色混合，factor 控制混合比例 (0=纯白, 1=原色)。"""
    return RGBColor(
        int((1 - factor) * 255 + factor * rgb[0]),
        int((1 - factor) * 255 + factor * rgb[1]),
        int((1 - factor) * 255 + factor * rgb[2]),
    )


def _shade(rgb, factor=0.7):
    """与黑色混合，factor 控制混合比例 (0=纯黑, 1=原色)。"""
    return RGBColor(
        int(factor * rgb[0]),
        int(factor * rgb[1]),
        int(factor * rgb[2]),
    )


def _luminance(rgb):
    """W3C 相对亮度 (sRGB → 线性 → 加权)。"""
    def _c(c):
        v = c / 255.0
        return v / 12.92 if v <= 0.03928 else ((v + 0.055) / 1.055) ** 2.4
    return 0.2126 * _c(rgb[0]) + 0.7152 * _c(rgb[1]) + 0.0722 * _c(rgb[2])


def contrast_text(rgb):
    """返回该底色上可读的文字色（白或黑），W3C 对比度 ≥ 4.5:1。"""
    L = _luminance(rgb)
    # 白色对比度: (1.0 + 0.05) / (L + 0.05)
    # 黑色对比度: (L + 0.05) / (0 + 0.05)
    white_contrast = 1.05 / (L + 0.05)
    black_contrast = (L + 0.05) / 0.05
    return WHITE if white_contrast > black_contrast else DARK


# ═══════════════════════════════════════════
# 主题色板（被 derive_palette() 覆盖）
# ═══════════════════════════════════════════

PALETTE_PRIMARY = TEAL       # 主题主色 → 卡片顶条、编号圆、重点标记
PALETTE_LIGHT = LIGHT_TEAL   # 主色浅变体 → 卡片背景
PALETTE_DARK = NAVY          # 暗色 → 章节页背景
PALETTE_WARM = LIGHT_GREEN   # 暖灰/第二卡片背景
PALETTE_ACCENT2 = GREEN_ACCENT  # 第二强调色


def derive_palette(primary_rgb, dark_rgb=None):
    """
    从主色自动生成一套协调色板（替代硬编码的 teal/green 默认值）。
    色调跟随主题，浅变体做卡片背景，暗色做章节页。
    """
    global PALETTE_PRIMARY, PALETTE_LIGHT, PALETTE_DARK, PALETTE_WARM, PALETTE_ACCENT2

    PALETTE_PRIMARY = primary_rgb
    PALETTE_LIGHT = _tint(primary_rgb, 0.12)       # 极浅 → 卡片背景
    PALETTE_DARK = dark_rgb if dark_rgb else _shade(primary_rgb, 0.15)
    # 暖灰：取主色的低饱和暖调偏移（HSV 色相不变，降饱和）
    PALETTE_WARM = RGBColor(245, 240, 235)          # 中性暖灰，与暖色主题均协调
    PALETTE_ACCENT2 = _shade(primary_rgb, 0.7)      # 深变体 → 第二强调


def init_design(json_path):
    """
    从 template_design.json 加载设计令牌，覆盖模块级默认值。
    如果文件不存在或解析失败，保持默认值不变。
    """
    global FONT_EN, BODY_SIZE, ACCENT_COLOR, TITLE_COLOR, TITLE_SIZE, BODY_COLOR, FOOTER_BARS, HEADER_DECOS

    if not json_path or not os.path.exists(json_path):
        return

    try:
        with open(json_path, 'r', encoding='utf-8') as f:
            d = json.load(f)
    except (json.JSONDecodeError, IOError):
        return

    # ── 字体 ──
    body = d.get('content_body', {})
    if body.get('font_name'):
        FONT_EN = body['font_name']
    if body.get('font_size_pt'):
        BODY_SIZE = Pt(body['font_size_pt'])
    if body.get('color'):
        BODY_COLOR = _hex_to_rgb(body['color'])

    # ── 内容标题 ──
    title = d.get('content_title', {})
    if title.get('font_size_pt'):
        TITLE_SIZE = Pt(title['font_size_pt'])
    if title.get('color'):
        TITLE_COLOR = _hex_to_rgb(title['color'])
        ACCENT_COLOR = TITLE_COLOR

    # ── 页脚 ──
    bars = d.get('footer_bars', [])
    if bars:
        FOOTER_BARS = []
        for b in bars:
            c = b.get('color', 'A5A5A5')
            FOOTER_BARS.append((b.get('y_in', 7.2), _hex_to_rgb(c), b.get('h_in', 0.1)))

    # ── 标题装饰 ──
    decos = d.get('header_decorations', [])
    if decos:
        HEADER_DECOS = []
        for dec in decos:
            HEADER_DECOS.append((
                dec.get('x_in', 0.3), dec.get('y_in', 0.25),
                dec.get('w_in', 0.15), dec.get('h_in', 0.45),
                _hex_to_rgb(dec.get('color', '007191'))
            ))

    # ── 主题色板推导 ──
    # 从设计令牌中提取主色和暗色，自动生成协调色板
    def _is_accent(c_hex):
        """判断是否为有彩度的强调色（排除灰、近黑、近白）。"""
        if c_hex in ('A5A5A5', 'CCCCCC', '999999', '888888', 'DDDDDD', 'EEEEEE'):
            return False
        rgb = _hex_to_rgb(c_hex)
        r, g, b = rgb[0], rgb[1], rgb[2]
        # 排除近黑 (max < 35)、近白 (min > 230)、低饱和灰色 (max-min < 20)
        if max(r, g, b) < 35 or min(r, g, b) > 230 or max(r, g, b) - min(r, g, b) < 20:
            return False
        return True

    primary = None
    dark_bg = None
    # 1. 主色：从 header_decorations 找有彩度的强调色
    for dec in (decos or []):
        c = dec.get('color', '')
        if _is_accent(c):
            primary = _hex_to_rgb(c)
            break
    # 2. 否则从 footer_bars 找
    if primary is None:
        for b in (d.get('footer_bars', []) or []):
            c = b.get('color', '')
            if _is_accent(c):
                primary = _hex_to_rgb(c)
                break
    # 3. 暗色：章节页背景 / 封面背景
    for key in ('section_bg_color', 'cover_bg_color'):
        if d.get(key):
            dark_bg = _hex_to_rgb(d[key])
            break
    # 4. 如果还是没有主色，用第一个 header 装饰的非全幅条
    if primary is None and decos:
        for dec in decos:
            if dec.get('w_in', 0) < 5:
                c = dec.get('color', '007191')
                primary = _hex_to_rgb(c)
                # 如果取到的还是近黑/灰，fallback 到默认
                if not _is_accent(c):
                    primary = None
                break

    if primary:
        derive_palette(primary, dark_bg)
        ACCENT_COLOR = primary
        if TITLE_COLOR == WHITE or TITLE_COLOR == _hex_to_rgb('FFFFFF'):
            TITLE_COLOR = primary

    # ── 如果存在 ppt-master 提取的 theme 数据，优先使用 ──
    theme = d.get('theme', {})
    if theme:
        # 使用 theme 中的权威 primary/dark 色覆盖启发式推导
        theme_primary = theme.get('palette_primary', '')
        theme_dark = theme.get('palette_dark', '')
        if theme_primary:
            primary_rgb = _hex_to_rgb(theme_primary)
            dark_rgb = _hex_to_rgb(theme_dark) if theme_dark else None
            derive_palette(primary_rgb, dark_rgb)
            ACCENT_COLOR = primary_rgb

        # 补充空白的封面报告人字体
        if 'cover_presenter' in d and not d['cover_presenter']:
            d['cover_presenter'] = {'font_name': FONT_EN, 'font_size_pt': 20, 'bold': False, 'color': BODY_COLOR}


def parse_color(val):
    """解析颜色值：RGBColor / 命名键 / '#RRGGBB'。失败返回 DARK。"""
    if isinstance(val, RGBColor):
        return val
    if val in COLOR_MAP:
        return COLOR_MAP[val]
    if val.startswith('#'):
        return _hex_to_rgb(val)
    return DARK


def _est_wrapped_lines(text_lines, font_sz, box_w_in):
    """估算多段文本在指定宽度内的实际折行总数。

    text_lines: 已按 \\n 分割的文本行列表
    font_sz:    EMU 字体大小 (如 Pt(14))
    box_w_in:   文本框宽度 (英寸)

    CJK 等宽 ~1 em，ASCII ~0.55 em。
    """
    font_in = font_sz / 12700 / 72  # pt → 英寸
    # PPT 文本框有内边距，有效宽度约为标称宽度的 92%
    chars_per_line = max(1, int(box_w_in / font_in * 0.92))
    total = 0
    for line in text_lines:
        if not line.strip():
            total += 1
            continue
        eff_len = 0
        for ch in line:
            if '一' <= ch <= '鿿' or '　' <= ch <= '〿' or '＀' <= ch <= '￯':
                eff_len += 1.0
            elif ord(ch) < 128:
                eff_len += 0.55
            else:
                eff_len += 1.0
        total += max(1, -(-int(eff_len) // chars_per_line))  # ceil division
    return max(1, total)


def _line_h(font_sz):
    """单行行高 (英寸)，含默认 120% 行距。"""
    return font_sz / 12700 / 72 * 1.25


def T(slide, l, t, w, h, text, sz=None, bold=False, color=None, align=PP_ALIGN.LEFT):
    """文本框，高度随文字自适应（不低于参数 h）。"""
    if sz is None:
        sz = BODY_SIZE
    if color is None:
        color = BODY_COLOR
    # 用 _est_wrapped_lines 估算折行（CJK≈1em, ASCII≈0.55em），比固定系数更准
    est_lines = _est_wrapped_lines([text], sz, w)
    actual_h = max(h, est_lines * _line_h(sz))
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(actual_h))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; r = p.add_run()
    r.text = text; r.font.name = FONT_EN; r.font.size = sz; r.font.bold = bold
    if color:
        r.font.color.rgb = parse_color(color) if isinstance(color, str) else color
    p.alignment = align
    return tb


def M(slide, l, t, w, h, lines, sz=None, color=None):
    """多行文本框，高度按实际折行估算（不低于参数 h）。

    逐字符加权计算折行数（CJK≈1em, ASCII≈0.55em），不再简单按段落数计算。
    """
    if sz is None: sz = BODY_SIZE
    if color is None: color = BODY_COLOR
    text_lines = [ln for ln in lines if ln.strip()] if lines else ['']
    if not text_lines:
        text_lines = ['']
    est_lines = _est_wrapped_lines(text_lines, sz, w)
    line_h_val = _line_h(sz)
    actual_h = max(h, est_lines * line_h_val + (len(text_lines) - 1) * 0.04)
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(actual_h))
    tf = tb.text_frame; tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        r = p.add_run(); r.text = line
        r.font.name = FONT_EN; r.font.size = sz
        r.font.color.rgb = parse_color(color) if isinstance(color, str) else color
        p.space_after = Pt(2)
    return tb


def _max_chars_fit(lines, font_sz, box_w, box_h):
    """估算在给定文本框内能容纳的最大字符数。返回可容纳的每条文本最大长度列表。"""
    font_in = font_sz / 12700 / 72
    chars_per_line = max(1, int(box_w / font_in))
    line_h_val = _line_h(font_sz)
    max_lines_total = max(1, int(box_h / line_h_val))
    # 为每个段落分配行数（平均分配，至少1行）
    n_paras = max(1, len(lines))
    lines_per_para = max(1, max_lines_total // n_paras)
    return [lines_per_para * chars_per_line] * n_paras


def truncate_lines(lines, font_sz, box_w, box_h, suffix='...'):
    """截断文本列表使其适配给定文本框。返回截断后的列表。"""
    if not lines:
        return lines
    limits = _max_chars_fit(lines, font_sz, box_w, box_h)
    result = []
    for i, text in enumerate(lines):
        limit = limits[min(i, len(limits) - 1)]
        if len(text) <= limit:
            result.append(text)
        else:
            # 截断，保留后缀空间
            cut = max(0, limit - len(suffix))
            result.append(text[:cut] + suffix)
    return result


def R(slide, l, t, w, h, fill_color, rounded=False, line_color=None):
    """填充矩形。"""
    st = 5 if rounded else 1
    rect = slide.shapes.add_shape(st, Inches(l), Inches(t), Inches(w), Inches(h))
    rect.fill.solid()
    rect.fill.fore_color.rgb = parse_color(fill_color) if isinstance(fill_color, str) else fill_color
    if line_color:
        rect.line.color.rgb = parse_color(line_color) if isinstance(line_color, str) else line_color
    else:
        rect.line.fill.background()
    return rect


def P(slide, path, l, t, w=None, max_h=None):
    """添加图片，保持宽高比。可指定 max_h 限制高度不超过。"""
    if w is None: w = IMG_W
    if not os.path.exists(path):
        print(f"  WARNING: {path} not found"); return None
    try:
        pic = slide.shapes.add_picture(path, Inches(l), Inches(t))
        ratio = pic.height / pic.width
        pic.width = Inches(w)
        pic.height = int(Inches(w) * ratio)
        if max_h and pic.height > Inches(max_h):
            pic.height = Inches(max_h)
            pic.width = int(Inches(max_h) / ratio)
        return pic
    except Exception as e:
        print(f"  ERROR adding {path}: {e}"); return None


def title_bar(slide, text):
    """
    内容页标题栏。使用设计令牌中的颜色和装饰。
    绘制顺序：背景条 → 装饰元素 → 标题文字（确保文字不被遮挡）。
    """
    # 1. 先画页眉装饰（背景条、色块、分隔线）
    if HEADER_DECOS:
        for x, y, w, h, c in HEADER_DECOS:
            r = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
            r.fill.solid(); r.fill.fore_color.rgb = c; r.line.fill.background()
    else:
        r = slide.shapes.add_shape(1, Inches(0.3), Inches(0.25), Inches(0.15), Inches(0.45))
        r.fill.solid(); r.fill.fore_color.rgb = ACCENT_COLOR; r.line.fill.background()
        l = slide.shapes.add_shape(1, Inches(0.6), Inches(0.68), Inches(12.1), Pt(1.5))
        l.fill.solid(); l.fill.fore_color.rgb = RGBColor(0xCC, 0xCC, 0xCC); l.line.fill.background()

    # 2. 再画标题文字（在最上层）
    T(slide, 0.5, 0.15, 12.3, 0.55, text, sz=TITLE_SIZE, bold=True, color=TITLE_COLOR)

    # 3. 页脚装饰条
    for y, c, h in FOOTER_BARS:
        b = slide.shapes.add_shape(1, Inches(0), Inches(y), Inches(13.33), Inches(h))
        b.fill.solid(); b.fill.fore_color.rgb = c; b.line.fill.background()


def white_bg(slide):
    """全幅白色背景。"""
    bg = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.33), Inches(7.5))
    bg.fill.solid(); bg.fill.fore_color.rgb = WHITE; bg.line.fill.background()


def blank_layout(prs):
    """查找空白布局。无空白布局时用占位符最少的布局，并清除残留占位符。"""
    for ly in prs.slide_layouts:
        if '空白' in ly.name or 'blank' in ly.name.lower():
            return ly
    # fallback: 选占位符最少的布局（避免 Title Slide 带入空框）
    best = min(prs.slide_layouts, key=lambda ly: len(ly.placeholders))
    return best


def make_result_slide(prs, title, body_lines, img_specs, figs_dir='.'):
    """
    标准结果内容页：左图右文，图片动态堆叠避免重叠/溢出。
    算法：预加载图片→计算实际高度→如溢出则等比限制→垂直堆叠。
    """
    slide = prs.slides.add_slide(blank_layout(prs))
    white_bg(slide)
    title_bar(slide, title)

    lines = [l for l in body_lines if l.strip()]
    n_lines = min(len(lines), 4)  # 最多 4 行
    img_x = 0.3
    area_top = 1.15
    area_bot = 6.85
    area_h = area_bot - area_top

    # ── 预加载图片 ──
    img_info = []
    for spec in img_specs:
        path = os.path.join(figs_dir, spec['file'])
        if os.path.exists(path):
            try:
                from PIL import Image
                im = Image.open(path)
                ratio = im.width / im.height
                img_info.append((path, ratio))
            except Exception:
                img_info.append((path, 1.5))
        else:
            img_info.append((path, 1.5))

    n_imgs = len(img_info)
    gap = 0.12  # 图片间距

    # ── 图片区布局（约束规则）──
    if n_imgs > 0:
        if n_imgs <= 2:
            # ≤2 图：垂直堆叠，每张不低于 2.0in 高
            img_w = IMG_W
            raw_heights = [img_w / r for _, r in img_info]
            if sum(raw_heights) + gap <= area_h:
                final_heights = raw_heights
                final_widths = [img_w] * n_imgs
            else:
                # 等分但每张不低于 2.0
                min_h = 2.0
                per_h = (area_h - (n_imgs - 1) * gap) / n_imgs
                final_heights = [max(min_h, per_h)] * n_imgs
                final_widths = [min(img_w, h * r) for h, (_, r) in zip(final_heights, img_info)]
            y = area_top
            for i, (path, _) in enumerate(img_info):
                P(slide, path, img_x, y, w=final_widths[i])
                y += final_heights[i] + gap
            img_end_y = y - gap
        else:
            # ≥3 图：2 列网格，每格不低于 1.3in
            n_cols = 2
            n_rows = (n_imgs + 1) // 2
            cell_w = (IMG_W - gap) / n_cols
            cell_h = min(2.2, (area_h - (n_rows - 1) * gap) / n_rows)
            cell_h = max(1.3, cell_h)
            for i, (path, ratio) in enumerate(img_info):
                col = i % 2; row = i // 2
                x = img_x + col * (cell_w + gap)
                y = area_top + row * (cell_h + gap)
                # 图片等比缩放到格子内，竖图(ratio<0.8)特殊处理
                if ratio < 0.8:
                    # 竖图：限高为格子的 80%，宽度自适应
                    pic_h = cell_h * 0.8
                    pic_w = min(cell_w, pic_h * ratio)
                else:
                    pic_w = cell_w
                    pic_h = min(cell_h, cell_w / ratio)
                pic = P(slide, path, x, y, w=pic_w)
                if pic:
                    pic.height = int(Inches(pic_h))
            img_end_y = area_top + n_rows * cell_h + (n_rows - 1) * gap

        # 文字区底对齐图片区底
        text_h = img_end_y - area_top
    else:
        text_h = area_h
        img_end_y = area_top

    # ── 文字区（加宽、自适应行距、不再重叠）──
    text_x = 6.8
    text_w = 6.2
    line_h = min(1.0, text_h / max(n_lines, 1))  # 行高自适应
    line_h = max(0.55, line_h)  # 每行至少 0.55in
    text_top = area_top + max(0, (text_h - n_lines * line_h) / 2)  # 垂直居中

    for i, line in enumerate(lines[:4]):
        y = text_top + i * line_h
        # 圆点标记
        dot = slide.shapes.add_shape(9, Inches(text_x - 0.25), Inches(y + 0.12),
                                      Inches(0.18), Inches(0.18))
        dot.fill.solid(); dot.fill.fore_color.rgb = ACCENT_COLOR; dot.line.fill.background()
        # 正文（行高比文本框小 0.05 留 margin）
        T(slide, text_x, y, text_w, line_h - 0.05, line)

    return slide
