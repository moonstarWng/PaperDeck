#!/usr/bin/env python3
"""
template_extractor.py — 从模板 PPTX 提取设计令牌 (design tokens)。
提取页眉/页脚、各级字体、颜色等，输出 template_design.json。

用法: python template_extractor.py <template.pptx> [output.json]
"""
from pptx import Presentation
from pptx.util import Pt
import json, sys, os


def _emu_to_pt(emu):
    """EMU → pt。"""
    if emu is None:
        return None
    return round(emu / 12700)


def _color_to_hex(rgb):
    """RGBColor → '#RRGGBB' 字符串。"""
    if rgb is None:
        return None
    try:
        return str(rgb)
    except:
        return None


def _find_shapes(slide, area=None, has_text=False, text_in=None):
    """
    查找幻灯片中的形状。
    area: 'header' (y<1.0), 'footer' (y>6.5), 'body' (1.0<y<6.5), None=all
    has_text: 仅返回有文本的形状
    text_in: 文本必须包含此字符串
    """
    results = []
    for shape in slide.shapes:
        y = shape.top / 914400
        h = shape.height / 914400
        w = shape.width / 914400

        # 过滤区域
        if area == 'header' and not (y < 1.0 and h < 1.5):
            continue
        if area == 'footer' and not (y > 6.5 and h < 0.5):
            continue
        if area == 'body' and not (1.0 < y < 6.5):
            continue

        # 过滤全幅背景（宽和高中至少一个超过全页的 90%）
        if w > 12 and h > 6.5:
            continue

        if has_text:
            if not shape.has_text_frame:
                continue
            if not shape.text_frame.text.strip():
                continue  # 跳过空文本框

        if text_in and shape.has_text_frame:
            full = shape.text_frame.text
            if text_in not in full:
                continue

        results.append(shape)
    return results


def _get_run_info(shape):
    """获取形状中第一个 run 的字体信息。"""
    if not shape.has_text_frame:
        return {}
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            f = run.font
            return {
                'font_name': f.name,
                'font_size_pt': _emu_to_pt(f.size),
                'bold': f.bold,
                'color': _color_to_hex(f.color.rgb) if f.color and f.color.rgb else None,
            }
    return {}


def _get_shape_style(shape):
    """获取形状的填充颜色。"""
    try:
        fill = shape.fill
        if fill.type is not None:
            return _color_to_hex(fill.fore_color.rgb)
    except:
        pass
    return None


def extract(template_path):
    """从模板 PPTX 提取设计令牌。"""
    prs = Presentation(template_path)
    total = len(prs.slides)
    if total < 5:
        raise ValueError(f"模板只有 {total} 页，至少需要 5 页 (cover/toc/sections/content/thanks)")

    # 识别页面类型
    from make_template import classify
    slide_types = {i: classify(s) for i, s in enumerate(prs.slides)}

    # 找到各类型的第一页
    indices = {}
    for i, t in slide_types.items():
        if t not in indices:
            indices[t] = i

    # 内容页检测：header + footer + body(含≤20pt文字)
    content_idx = None
    for i in range(total):
        s = prs.slides[i]
        if slide_types.get(i, '') in ('TOC', 'THANKS'):
            continue
        headers = _find_shapes(s, area='header', has_text=True)
        footers = _find_shapes(s, area='footer')
        bodies = _find_shapes(s, area='body', has_text=True)
        if not (headers and footers and len(bodies) >= 2):
            continue
        # 内容页特征：header title ≤28pt; body 有 ≤20pt 正文
        hdr_too_big = False
        for sh in headers:
            info = _get_run_info(sh)
            if info.get('font_size_pt') and info['font_size_pt'] > 28:
                hdr_too_big = True
        if hdr_too_big:
            continue
        has_body_text = False
        for sh in bodies:
            info = _get_run_info(sh)
            if info.get('font_size_pt') and info['font_size_pt'] <= 20:
                has_body_text = True
                break
        if has_body_text:
            content_idx = i
            break
    if content_idx is None:
        content_idx = total - 2

    cover_idx = indices.get('COVER', 0)
    section_indices = [i for i, t in slide_types.items() if t.startswith('SECTION_')]
    section_idx = section_indices[0] if section_indices else 2
    toc_idx = indices.get('TOC', 1)
    thanks_idx = indices.get('THANKS', total - 1)

    # ═══ 提取 ═══
    design = {'meta': {'source': os.path.basename(template_path), 'total_slides': total}}

    # ── 内容页 ──
    cs = prs.slides[content_idx]
    title_shapes = _find_shapes(cs, area='header', has_text=True)
    body_shapes = _find_shapes(cs, area='body', has_text=True)
    footer_shapes = _find_shapes(cs, area='footer')
    header_decos = _find_shapes(cs, area='header', has_text=False)

    # 标题字体（header 区域第一个有文本的形状）
    design['content_title'] = _get_run_info(title_shapes[0]) if title_shapes else {}
    # 正文字体（body 区域最大的文本框）
    if body_shapes:
        largest = max(body_shapes, key=lambda s: s.height)
        design['content_body'] = _get_run_info(largest)
    else:
        design['content_body'] = {}

    # 页脚颜色
    footer_colors = []
    for sh in footer_shapes:
        c = _get_shape_style(sh)
        if c:
            footer_colors.append({'color': c, 'y_in': round(sh.top / 914400, 2), 'h_in': round(sh.height / 914400, 3)})
    design['footer_bars'] = footer_colors

    # 标题装饰（header 区域非文本形状）
    for sh in header_decos:
        c = _get_shape_style(sh)
        if c:
            design.setdefault('header_decorations', []).append({
                'color': c, 'x_in': round(sh.left / 914400, 2), 'y_in': round(sh.top / 914400, 2),
                'w_in': round(sh.width / 914400, 3), 'h_in': round(sh.height / 914400, 3),
            })

    # ── 封面 ──
    cover = prs.slides[cover_idx]
    # 封面标题：大号、粗体文本
    cover_texts = [s for s in cover.shapes if s.has_text_frame]
    cover_fonts = []
    for s in cover_texts:
        info = _get_run_info(s)
        if info.get('font_size_pt', 0) and info['font_size_pt'] > 20:
            cover_fonts.append(info)
    design['cover_title'] = cover_fonts[0] if cover_fonts else {}
    design['cover_presenter'] = cover_fonts[-1] if len(cover_fonts) > 1 else {}

    # 封面背景色
    bg_shapes = [s for s in cover.shapes if s.width / 914400 > 12 and s.height / 914400 > 7]
    if bg_shapes:
        c = _get_shape_style(bg_shapes[0])
        if c:
            design['cover_bg_color'] = c

    # ── 章节页 ──
    sec = prs.slides[section_idx]
    sec_texts = [s for s in sec.shapes if s.has_text_frame]
    for s in sec_texts:
        txt = s.text_frame.text.strip()
        info = _get_run_info(s)
        if txt.isdigit() and len(txt) <= 2:
            design['section_number'] = info
        elif len(txt) > 1:
            design['section_title'] = info
    bg_shapes = [s for s in sec.shapes if s.width / 914400 > 12 and s.height / 914400 > 7]
    if bg_shapes:
        c = _get_shape_style(bg_shapes[0])
        if c:
            design['section_bg_color'] = c

    # ── 目录页 ──
    toc = prs.slides[toc_idx]
    toc_items = [s for s in toc.shapes if s.has_text_frame and s.top / 914400 > 1.5]
    toc_fonts = []
    for s in toc_items:
        info = _get_run_info(s)
        txt = s.text_frame.text.strip()
        if txt.isdigit():
            design['toc_number'] = info
        elif len(txt) > 1:
            design['toc_item'] = info

    # ── 致谢页 ──
    thanks = prs.slides[thanks_idx]
    thanks_texts = [s for s in thanks.shapes if s.has_text_frame]
    for s in thanks_texts:
        info = _get_run_info(s)
        if info.get('font_size_pt', 0) and info['font_size_pt'] > 30:
            design['thanks_title'] = info
            break

    # ── 默认正文字体（fallback）──
    if design.get('content_body', {}).get('font_name'):
        design['default_font'] = design['content_body']['font_name']

    return design


def main():
    if len(sys.argv) < 2:
        print(f"Usage: python {sys.argv[0]} <template.pptx> [output.json]")
        sys.exit(1)

    src = sys.argv[1]
    dst = sys.argv[2] if len(sys.argv) > 2 else src.replace('.pptx', '_design.json')

    if not os.path.exists(src):
        print(f"ERROR: {src} not found")
        sys.exit(1)

    design = extract(src)
    with open(dst, 'w', encoding='utf-8') as f:
        json.dump(design, f, indent=2, ensure_ascii=False)

    print(f"提取完成 → {dst}")
    print(f"  {len(design)} 个设计令牌")
    for k, v in design.items():
        if k != 'meta' and k != 'footer_bars' and k != 'header_decorations':
            print(f"  {k}: {v}")


if __name__ == '__main__':
    main()
