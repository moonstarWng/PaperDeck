#!/usr/bin/env python3
"""
llm_template_extract.py — LLM 驱动的模板骨架提取。

与 make_template.py（纯规则）互补：
  - make_template.py: 快速、确定性，但 Logo/图标阈值硬编码
  - llm_template_extract.py: 自适应不同模板设计，需要 LLM API

工作流程:
  1. dump_slide_info()    — 把每页结构 dump 为 JSON
  2. classify_slides()    — 调 LLM 批量分类 (1 次调用)
  3. identify_elements()  — 调 LLM 识别 CONTENT 页元素 (N 次，可并行)
  4. apply_placeholders() — 替换内容为占位文本/图片

用法:
  python llm_template_extract.py <input.pptx> [output.pptx]
"""

import sys, os, json, re
from pptx import Presentation
from pptx.util import Inches, Pt


# ═══════════════════════════════════════════
# Step 1: 幻灯片结构化提取
# ═══════════════════════════════════════════

def _safe_fill_color(shape):
    """安全获取形状填充颜色，返回 '#RRGGBB' 或 None。"""
    try:
        fill = shape.fill
        if fill.type is not None:
            return str(fill.fore_color.rgb)
    except Exception:
        pass
    return None


def _safe_font_info(shape):
    """从形状的第一个文本 run 中提取字体信息。"""
    if not shape.has_text_frame:
        return None
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            f = run.font
            info = {}
            if f.size:
                info['size_pt'] = round(f.size / 12700)
            if f.bold is not None:
                info['bold'] = f.bold
            try:
                if f.color.rgb:
                    info['color'] = str(f.color.rgb)
            except Exception:
                pass
            if f.name:
                info['name'] = f.name
            return info if info else None
    return None


def dump_slide_info(slide):
    """
    将单个幻灯片的所有形状提取为 LLM 可读的结构化 JSON。
    保留信息: 类型、名称、位置、尺寸、填充色、文本、字体。
    """
    shapes_info = []
    for shape in slide.shapes:
        stype = str(shape.shape_type)
        info = {
            'name': shape.name,
            'type': stype,
            'pos': {
                'l': round(shape.left / 914400, 2),
                't': round(shape.top / 914400, 2),
            },
            'size': {
                'w': round(shape.width / 914400, 2),
                'h': round(shape.height / 914400, 2),
            },
        }

        # 填充色
        fill = _safe_fill_color(shape)
        if fill:
            info['fill'] = fill

        # 文本 + 字体
        if shape.has_text_frame:
            text = shape.text_frame.text.strip()
            if text:
                info['text'] = text[:200]  # 截断过长文本
            font = _safe_font_info(shape)
            if font:
                info['font'] = font

        shapes_info.append(info)

    return {
        'shapes': shapes_info,
        'shape_count': len(shapes_info),
    }


def dump_all_slides(prs):
    """提取所有幻灯片的结构信息。"""
    slides_data = []
    for i, slide in enumerate(prs.slides):
        info = dump_slide_info(slide)
        info['slide_index'] = i
        slides_data.append(info)
    return slides_data


# ═══════════════════════════════════════════
# Step 2: LLM 幻灯片分类
# ═══════════════════════════════════════════

CLASSIFICATION_SYSTEM = """你是一个 PPT 结构分析专家。你的任务是根据幻灯片的结构特征，
将每页分类为以下 5 种类型之一。

## 分类定义

### COVER (封面)
- 大号标题文字 (>28pt)、汇报人/机构/日期信息
- 通常有大面积纯色背景装饰（左侧竖条或全幅背景）
- 文本框数量 ≤ 4 个
- 典型文本：「论文标题」「汇报人」「XX大学」「202X年X月」

### TOC (目录)
- 列表式结构，包含 "目录"/"CONTENTS"/"Contents" 或编号 01/02/03
- GROUP 形状较多（代表复杂的排版组合）
- 文本框数量 > 4 个
- 典型文本：「目录」「CONTENTS」「01」「02」

### SECTION (章节分隔页)
- 单一章节编号 (01-09) + 简短章节标题
- 大面积深色背景（填充色为深色）
- 文本框数量 ≤ 3 个
- 典型文本：「01 研究背景」「02 Methods」

### CONTENT (内容页)
- 有标题栏区域 (通常在页面上方，字体 >16pt)
- 有正文文本区域 (多段文字，每段独立文本框)
- 常有图片区域
- 有页脚装饰色条
- 文本框数量 ≥ 3 个，且不是纯列表结构

### THANKS (尾页/致谢)
- 极简布局，仅 1-2 个文本框
- 大字 (>36pt) 居中
- 典型文本：「THANKS」「谢谢」「感谢聆听」「Q&A」

## 输出要求
返回一个 JSON 对象，格式为:
{
  "classifications": {
    "0": "COVER",
    "1": "TOC",
    ...
  }
}

键是 slide_index，值是分类名 (COVER/TOC/SECTION/CONTENT/THANKS)。
不要包裹在 ```json``` 中，不要有任何解释文字。"""


def build_classification_prompt(slides_data):
    """构建分类提示词。"""
    # 发送精简版结构（去掉过细的字体信息，减少 token）
    compact = []
    for sd in slides_data:
        shapes_summary = []
        for s in sd['shapes']:
            summary = {
                'type': s['type'],
                'size': f"{s['size']['w']:.1f}x{s['size']['h']:.1f}in",
            }
            if 'fill' in s:
                summary['fill'] = s['fill']
            text = s.get('text', '')
            if text:
                summary['text'] = text[:120]
            if s.get('font'):
                summary['font_size'] = s['font'].get('size_pt', '?')
            shapes_summary.append(summary)
        compact.append({
            'slide_index': sd['slide_index'],
            'shape_count': sd['shape_count'],
            'shapes': shapes_summary,
        })
    return f"请对以下 {len(compact)} 页幻灯片进行分类:\n\n{json.dumps(compact, indent=2, ensure_ascii=False)}"


def parse_classification_response(response_text):
    """解析 LLM 的分类响应。"""
    # 清理可能的 markdown 包裹和思考文本
    text = response_text.strip()
    # 提取 JSON 块
    m = re.search(r'\{[\s\S]*"classifications"[\s\S]*\}', text)
    if m:
        text = m.group(0)
    try:
        data = json.loads(text)
        return {int(k): v for k, v in data.get('classifications', {}).items()}
    except json.JSONDecodeError:
        # 尝试修复常见问题 (trailing commas, single quotes)
        text = re.sub(r',\s*}', '}', text)
        text = re.sub(r',\s*]', ']', text)
        try:
            data = json.loads(text)
            return {int(k): v for k, v in data.get('classifications', {}).items()}
        except json.JSONDecodeError:
            raise ValueError(f"无法解析 LLM 响应: {response_text[:500]}")


# ═══════════════════════════════════════════
# Step 3: LLM 内容页元素识别
# ═══════════════════════════════════════════

ELEMENT_SYSTEM = """你是一个 PPT 版式分析专家。你的任务是识别内容页幻灯片中每个形状的"角色"。

## 角色定义

### title (标题区)
- 字体较大 (>16pt)、通常加粗
- 位于页面上方 (y < 1.5in)
- 文字简短（一行）
- 可能位于深色标题条背景上

### body (正文区)
- 字体较小 (≤16pt)
- 位于页面中部 (y 在 1.0-6.5in 之间)
- 多段文字，每段独立文本框
- 前面可能有小圆点标记

### image (图片区)
- PICTURE 类型形状
- 或较大的 AUTO_SHAPE/TEXT_BOX 位置（可能是图片占位区）
- 本身不含文本或含极短文本（如图片编号）

### decoration (装饰元素)
- 无文字的小形状 (线条、色块、圆点、分隔线)
- 位于页眉或页脚区域
- 填充色为纯色的小矩形或圆形
- 全幅背景矩形 (w > 10in or h > 6in)

## 输出要求
返回一个 JSON 对象:
{
  "elements": {
    "0": "decoration",    // shape 的序号 → 角色
    "1": "title",
    "2": "body",
    "3": "image",
    ...
  }
}
键是形状在 shapes 数组中的索引，值是角色名。
不要包裹在 ```json``` 中，不要有任何解释文字。"""


def build_element_prompt(slide_data):
    """为单个内容页构建元素识别提示词。"""
    # 为每个形状编号
    shapes_with_id = []
    for i, s in enumerate(slide_data['shapes']):
        entry = {
            'id': i,
            'type': s['type'],
            'name': s['name'],
            'pos': s['pos'],
            'size': s['size'],
        }
        if 'fill' in s:
            entry['fill'] = s['fill']
        if 'text' in s:
            entry['text'] = s['text'][:150]
        if 'font' in s:
            entry['font'] = s['font']
        shapes_with_id.append(entry)

    return f"请分析以下内容页的 {len(shapes_with_id)} 个形状:\n\n{json.dumps(shapes_with_id, indent=2, ensure_ascii=False)}"


def parse_element_response(response_text):
    """解析元素识别响应。"""
    text = response_text.strip()
    m = re.search(r'\{[\s\S]*"elements"[\s\S]*\}', text)
    if m:
        text = m.group(0)
    try:
        data = json.loads(text)
        return {int(k): v for k, v in data.get('elements', {}).items()}
    except json.JSONDecodeError:
        text = re.sub(r',\s*}', '}', text)
        text = re.sub(r',\s*]', ']', text)
        try:
            data = json.loads(text)
            return {int(k): v for k, v in data.get('elements', {}).items()}
        except json.JSONDecodeError:
            raise ValueError(f"无法解析元素响应: {response_text[:500]}")


# ═══════════════════════════════════════════
# Step 4: 应用占位符替换
# ═══════════════════════════════════════════

PLACEHOLDER_COLOR = (0x99, 0x99, 0x99)  # 灰色占位文字/图片背景


def apply_placeholders(prs, slide_index, element_map):
    """
    根据元素角色映射，将内容页的文字和图片替换为占位符。
    - title → "此处填充标题"
    - body → "此处填充文本"
    - image → 替换为灰色占位矩形 "此处填充图片"
    - decoration → 保留不动
    """
    from pptx.dml.color import RGBColor

    slide = prs.slides[slide_index]
    gray = RGBColor(*PLACEHOLDER_COLOR)
    shapes_list = list(slide.shapes)

    for idx, role in element_map.items():
        if idx >= len(shapes_list):
            continue
        shape = shapes_list[idx]

        if role == 'title':
            if shape.has_text_frame:
                # 只改文字，保留原始字体
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if run.text.strip():
                            run.text = '此处填充标题'

        elif role == 'body':
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if run.text.strip():
                            run.text = '此处填充文本'

        elif role == 'image':
            is_picture = 'PICTURE' in str(shape.shape_type)
            if is_picture:
                # 删除图片，在同位置放一个灰色占位矩形
                l, t = shape.left, shape.top
                w, h = shape.width, shape.height
                shape._element.getparent().remove(shape._element)
                rect = slide.shapes.add_shape(1, l, t, w, h)  # AUTO_SHAPE rectangle
                rect.fill.solid()
                rect.fill.fore_color.rgb = RGBColor(0xE0, 0xE0, 0xE0)
                rect.line.fill.background()
                tf = rect.text_frame
                tf.word_wrap = True
                p = tf.paragraphs[0]
                p.alignment = 1  # PP_ALIGN.CENTER
                r = p.add_run()
                r.text = '此处填充图片'
                r.font.size = Pt(14)
                r.font.color.rgb = gray
            else:
                # 非图片但被标记为 image 区的形状：如果无文字则保留，有文字则替换
                if shape.has_text_frame:
                    tf = shape.text_frame
                    tf.clear()
                    p = tf.paragraphs[0]
                    r = p.add_run()
                    r.text = '此处填充图片'
                    r.font.size = Pt(14)
                    r.font.color.rgb = gray

        # decoration → 不做任何修改


# ═══════════════════════════════════════════
# 持久化分析缓存
# ═══════════════════════════════════════════

def _save_analysis_cache(cache_path, slides_data, classifications, element_maps, prs):
    """将模板分析结果持久化为 JSON 缓存文件。"""
    import hashlib, datetime
    slides_cache = []
    for i, sd in enumerate(slides_data):
        cat = classifications.get(i, 'UNKNOWN')
        emap = element_maps.get(i, {})
        shapes_out = []
        for j, s in enumerate(sd['shapes']):
            entry = {
                'name': s['name'], 'type': s['type'],
                'pos': s['pos'], 'size': s['size'],
                'role': emap.get(j, 'unknown'),
            }
            if 'fill' in s: entry['fill'] = s['fill']
            if 'text' in s: entry['text'] = s['text'][:100]
            if 'font' in s: entry['font'] = s['font']
            shapes_out.append(entry)
        slides_cache.append({
            'index': i, 'classification': cat,
            'shapes': shapes_out,
        })
    # 提取设计令牌（失败不影响主流程）
    design = {}
    try:
        import tempfile
        tmp = tempfile.NamedTemporaryFile(suffix='.pptx', delete=False)
        tmp.close()
        prs.save(tmp.name)
        from template_extractor import extract as extract_design
        design = extract_design(tmp.name)
        os.unlink(tmp.name)
    except Exception:
        pass

    cache = {
        'source': os.path.basename(cache_path.replace('_analysis.json', '.pptx')),
        'analyzed_at': datetime.datetime.now().isoformat(),
        'total_slides': len(slides_cache),
        'slides': slides_cache,
        'design_tokens': design,
    }
    with open(cache_path, 'w', encoding='utf-8') as f:
        json.dump(cache, f, indent=2, ensure_ascii=False)


def _extract_design_inline(prs):
    """内联提取设计令牌（不依赖文件路径）。"""
    from template_extractor import extract as _ext
    import tempfile
    tmp = tempfile.NamedTemporaryFile(suffix='.pptx', delete=False)
    tmp.close()
    try:
        prs.save(tmp.name)
        return _ext(tmp.name)
    finally:
        try: os.unlink(tmp.name)
        except: pass


def _process_dir(pptx_path):
    """获取过程文件目录（论文同级 process/）。"""
    d = os.path.join(os.path.dirname(os.path.abspath(pptx_path)), 'process')
    os.makedirs(d, exist_ok=True)
    return d

def load_template_cache(pptx_path):
    """加载模板分析缓存。返回 None 表示缓存不存在或已过期。"""
    cache_path = os.path.join(_process_dir(pptx_path), 'analysis.json')
    if not os.path.exists(cache_path):
        return None
    pptx_mtime = os.path.getmtime(pptx_path)
    cache_mtime = os.path.getmtime(cache_path)
    if cache_mtime < pptx_mtime:
        return None  # PPTX 已更新，缓存过期
    try:
        with open(cache_path, 'r', encoding='utf-8') as f:
            return json.load(f)
    except Exception:
        return None


def analyze_and_cache(pptx_path):
    """分析 PPTX 并持久化缓存（规则分类，无需 LLM）。返回缓存 dict。"""
    prs = Presentation(pptx_path)
    slides_data = dump_all_slides(prs)

    sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))
    from make_template import classify

    classifications = {}
    for i, slide in enumerate(prs.slides):
        cat = classify(slide)
        if cat.startswith('SECTION_'): classifications[i] = 'SECTION'
        elif cat in ('COVER', 'TOC', 'THANKS'): classifications[i] = cat
        else: classifications[i] = 'CONTENT'

    # 内容页：按规则分配角色（标题/正文/图片/装饰）
    element_maps = {}
    for i in classifications:
        if classifications[i] != 'CONTENT':
            continue
        slide = prs.slides[i]
        emap = {}
        for j, shape in enumerate(slide.shapes):
            role = 'decoration'
            if shape.has_text_frame and shape.text_frame.text.strip():
                y = shape.top / 914400
                font = _safe_font_info(shape)
                size = font.get('size_pt', 0) if font else 0
                if y < 1.5 and size > 16:
                    role = 'title'
                elif y < 7.0:
                    role = 'body'
            elif 'PICTURE' in str(shape.shape_type):
                role = 'image'
            emap[j] = role
        element_maps[i] = emap

    cache_path = os.path.join(_process_dir(pptx_path), 'analysis.json')
    _save_analysis_cache(cache_path, slides_data, classifications, element_maps, prs)
    return load_template_cache(pptx_path)


# ═══════════════════════════════════════════
# 主入口
# ═══════════════════════════════════════════

def extract_template_llm(input_path, output_path, base_url, api_key, model,
                         on_progress=None):
    """
    LLM 驱动的模板骨架提取。

    参数:
      input_path: 完整 PPTX 文件路径
      output_path: 输出模板骨架路径
      base_url, api_key, model: LLM API 参数
      on_progress: 进度回调 (message: str) -> None

    返回:
      dict: 包含分类结果和元素识别结果
    """
    import requests

    def log(msg):
        if on_progress:
            on_progress(msg)
        print(f"  {msg}")

    log(f"读取 {input_path}...")
    prs = Presentation(input_path)
    total = len(prs.slides)
    log(f"共 {total} 页")

    # ── Step 1: dump ──
    log("提取幻灯片结构...")
    slides_data = dump_all_slides(prs)

    # ── Step 2: 分类（规则优先，保证封面/目录/章节/致谢不丢）──
    log("规则分类...")
    sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))
    from make_template import classify
    classifications = {}
    for i, slide in enumerate(prs.slides):
        cat = classify(slide)
        # 映射到简化分类
        if cat.startswith('SECTION_'):
            classifications[i] = 'SECTION'
        elif cat in ('COVER', 'TOC', 'THANKS'):
            classifications[i] = cat
        else:
            classifications[i] = 'CONTENT'
    log(f"分类结果: {json.dumps(classifications, ensure_ascii=False)}")

    # ── Step 3: 规则替换（不调 LLM，全部文本和图片直接替换）──
    content_slides = [i for i, t in classifications.items() if t == 'CONTENT']
    log(f"内容页: {len(content_slides)} 页，规则替换...")

    from pptx.dml.color import RGBColor
    gray = RGBColor(*PLACEHOLDER_COLOR)
    for idx in content_slides:
        slide = prs.slides[idx]
        for shape in list(slide.shapes):
            if 'PICTURE' in str(shape.shape_type):
                l, t, w, h = shape.left, shape.top, shape.width, shape.height
                shape._element.getparent().remove(shape._element)
                rect = slide.shapes.add_shape(1, l, t, w, h)
                rect.fill.solid(); rect.fill.fore_color.rgb = RGBColor(0xE0, 0xE0, 0xE0)
                rect.line.fill.background()
                tf = rect.text_frame; tf.word_wrap = True
                p = tf.paragraphs[0]; p.alignment = 1
                r = p.add_run(); r.text = '此处填充图片'; r.font.size = Pt(14); r.font.color.rgb = gray
            elif shape.has_text_frame and shape.text_frame.text.strip():
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if run.text.strip():
                            run.text = '此处填充文本'

    # 删除非内容页的所有图片
    for i in range(total):
        if i not in content_slides:
            for s in list(prs.slides[i].shapes):
                if 'PICTURE' in str(s.shape_type):
                    s._element.getparent().remove(s._element)

    # 通用化非内容页文字
    log("通用化模板文字...")
    from make_template import generalize_text
    generalize_text(prs)

    # ── 持久化分析缓存 ──
    cache_path = os.path.join(_process_dir(input_path), 'analysis.json')
    _save_analysis_cache(cache_path, slides_data, classifications, all_element_maps, prs)

    # ── 保存 ──
    log(f"保存 → {output_path}")
    os.makedirs(os.path.dirname(output_path) or '.', exist_ok=True)
    prs.save(output_path)

    # ── 返回结果 ──
    return {
        'classifications': classifications,
        'element_maps': all_element_maps,
        'total_slides': total,
        'content_slides': content_slides,
    }


# ═══════════════════════════════════════════
# CLI 入口
# ═══════════════════════════════════════════

def main():
    if len(sys.argv) < 2:
        print('Usage: python llm_template_extract.py <input.pptx> [output.pptx]')
        print('Environment variables:')
        print('  LLM_BASE_URL  — API endpoint (default: http://localhost:8000)')
        print('  LLM_API_KEY   — API key (default: none)')
        print('  LLM_MODEL     — Model name (default: gpt-4o)')
        sys.exit(1)

    src = sys.argv[1]
    dst = sys.argv[2] if len(sys.argv) > 2 else src.replace('.pptx', '_llm_template.pptx')

    base_url = os.environ.get('LLM_BASE_URL', 'http://localhost:8000')
    api_key = os.environ.get('LLM_API_KEY', '')
    model = os.environ.get('LLM_MODEL', 'gpt-4o')

    result = extract_template_llm(src, dst, base_url, api_key, model)
    print(f"\n完成: {dst}")
    print(f"分类: {json.dumps(result['classifications'], ensure_ascii=False)}")
    print(f"内容页数: {len(result['content_slides'])}")


if __name__ == '__main__':
    main()
