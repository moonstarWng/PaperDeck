"""
ppt_slides.py — 参数化的特殊幻灯片构建器。
每个函数接收一个从 slide-content.json 解析出的 data 字典，
构建特定类型的幻灯片（作者团队、课题背景、结果总结、讨论等）。
所有函数遵循相同的签名: (prs, data) -> slide
"""
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os, sys
sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))
from ppt_layout import (
    T, M, R, P, title_bar, white_bg, blank_layout,
    TEAL, DARK, WHITE, NAVY,  # 保留旧常量供 fallback
    COLOR_MAP, parse_color, contrast_text,
    BODY_SIZE, FONT_EN, FONT_CN,
    PALETTE_PRIMARY, PALETTE_LIGHT, PALETTE_DARK, PALETTE_WARM, PALETTE_ACCENT2,
    _est_wrapped_lines, _line_h, truncate_lines,
)



def build_author_slide(prs, data):
    """
    构建"作者团队与研究背景"页。
    布局：上方左右两张卡片（期刊信息 + 研究机构），下方通栏横幅（通讯作者 + 前期基础）。

    data 结构:
      journal: {name, if, doi, date}
      institutions: [字符串列表]
      authors: 通讯作者字符串（用 | 分隔）
      prior_work: [前期研究工作列表]
    """
    slide = prs.slides.add_slide(blank_layout(prs))
    white_bg(slide)
    title_bar(slide, '1 作者团队与研究背景')

    journal = data.get('journal', {})
    institutions = data.get('institutions', [])
    authors = data.get('authors', '')
    prior_work = data.get('prior_work', [])

    # ── 左上卡片：期刊信息 ──
    journal_lines = 3  # 名称+IF, 日期, DOI
    jcard_h = max(2.2, 0.8 + journal_lines * 0.5)
    R(slide, 0.4, 1.1, 5.8, jcard_h, PALETTE_LIGHT, rounded=True)
    R(slide, 0.4, 1.1, 5.8, 0.5, PALETTE_PRIMARY)
    T(slide, 0.6, 1.15, 5.4, 0.4, '  期刊信息', sz=Pt(18), bold=True, color=WHITE)
    # 期刊信息行：避免空值产生多余分隔符
    journal_info_lines = []
    jname = journal.get('name', '')
    # 第1行：期刊全名 + IF
    if jname:
        if journal.get('if'):
            journal_info_lines.append(f"{jname} (IF ~{journal['if']})")
        else:
            journal_info_lines.append(jname)
    # 第2行：子刊类型 + 日期（仅当有实际内容时）
    date_part = journal.get('date', '')
    nature_part = 'Nature 子刊' if 'Nature' in jname else ''
    if nature_part or date_part:
        journal_info_lines.append(f"{nature_part} | {date_part}" if nature_part and date_part else (nature_part or date_part))
    # 第3行：DOI
    if journal.get('doi'):
        journal_info_lines.append(f"DOI: {journal['doi']}")
    M(slide, 0.7, 1.8, 5.2, jcard_h - 0.8, journal_info_lines, sz=BODY_SIZE, color=DARK)

    # ── 右上卡片：研究机构 ──
    inst_lines = max(len(institutions), 2)
    icard_h = max(2.2, 0.8 + inst_lines * 0.5)
    R(slide, 6.6, 1.1, 6.3, icard_h, PALETTE_WARM, rounded=True)
    R(slide, 6.6, 1.1, 6.3, 0.5, PALETTE_ACCENT2)
    T(slide, 6.8, 1.15, 5.9, 0.4, '  研究机构', sz=Pt(18), bold=True, color=WHITE)
    M(slide, 6.9, 1.8, 5.8, icard_h - 0.8, institutions, sz=BODY_SIZE, color=DARK)

    # ── 底部通栏：通讯作者 + 前期基础 ──
    prior_lines = len(prior_work) if prior_work else 0
    author_bot_h = max(2.2, 1.5 + prior_lines * 0.55)
    author_bot_h = min(5.5, author_bot_h)
    R(slide, 0.4, 3.9, 12.5, author_bot_h, PALETTE_LIGHT, rounded=True)
    T(slide, 0.7, 4.0, 5.0, 0.4, '  通讯作者', sz=Pt(18), bold=True, color=PALETTE_PRIMARY)
    T(slide, 0.7, 4.5, 11.8, 0.4, authors, sz=BODY_SIZE, bold=True, color=PALETTE_PRIMARY)
    T(slide, 0.7, 5.1, 5.0, 0.4, '  课题组前期研究基础', sz=BODY_SIZE, bold=True, color=PALETTE_PRIMARY)
    if prior_work:
        M(slide, 0.7, 5.5, 11.8, author_bot_h - 1.3, prior_work, sz=BODY_SIZE, color=DARK)

    return slide


def build_background_slide(prs, data):
    """
    构建"课题背景"页。
    布局：3 张横向排列的卡片（每张含彩色标题条） + 底部假说/实验横幅。

    data 结构:
      cards: [{title, body, color}, ...]  — 2-4 张卡片
      hypothesis: 核心假说文本
      experiment: 实验设计文本
    """
    slide = prs.slides.add_slide(blank_layout(prs))
    white_bg(slide)
    title_bar(slide, '2 课题背景：从临床问题到科学机制')

    cards = data.get('cards', [])
    n_cards = len(cards)
    margin = 0.45
    gap_x = 0.27
    card_w = (13.33 - margin * 2 - gap_x * (n_cards - 1)) / max(n_cards, 1)
    card_w = min(card_w, 8.0)
    # 用原始 \n 分段估算折行高度（不硬拆行，PPT 自然换行）
    card_body_w = card_w - 0.35
    max_wrapped = 1
    for card in cards:
        body = card.get("body","").strip()
        bl = [l for l in body.split(chr(10)) if l.strip()]
        if bl:
            n = _est_wrapped_lines(bl, BODY_SIZE, card_body_w)
            max_wrapped = max(max_wrapped, n)
    # 段落内折行的真实行高
    body_line_h = _line_h(BODY_SIZE)
    card_h_ideal = max(2.8, 1.1 + max_wrapped * body_line_h)
    card_h = min(5.4, card_h_ideal)  # 上限：留 2.0in 给底部假说/实验横幅
    # 正文溢出时自动缩小字号
    avail_body_h = card_h - 0.9
    body_need_h = max_wrapped * body_line_h
    body_sz = BODY_SIZE
    if body_need_h > avail_body_h * 1.05:
        scale = avail_body_h / body_need_h
        body_sz = Pt(max(10, int(BODY_SIZE / 12700 * scale)))
    for i, card in enumerate(cards):
        cx = margin + i * (card_w + gap_x)
        color = parse_color(card.get("color", "teal"))
        bl = [l for l in card.get("body","").strip().split(chr(10)) if l.strip()]
        R(slide, cx, 1.15, card_w, card_h, PALETTE_LIGHT, rounded=True)
        R(slide, cx, 1.15, card_w, 0.55, color)
        title_sz = Pt(18) if n_cards >= 3 else Pt(20)
        T(slide, cx + 0.15, 1.2, card_w - 0.35, 0.45, card["title"], sz=title_sz, bold=True, color=WHITE)
        # 截断过长卡片文本
        bl_trunc = truncate_lines(bl, body_sz, card_w - 0.35, card_h - 0.9)
        M(slide, cx + 0.15, 1.85, card_w - 0.35, card_h - 0.9, bl_trunc, sz=body_sz, color=DARK)

    # 底部假说/实验横幅（位置跟随卡片高度动态调整）
    hypothesis = data.get('hypothesis', '')
    experiment = data.get('experiment', '')
    if hypothesis or experiment:
        banner_y = 1.15 + card_h + 0.2  # 卡片底部 + 间距
        banner_h = 1.8
        # 确保不超出幻灯片底部 (7.5 - 页脚 ~0.3 = 7.2)
        if banner_y + banner_h > 7.0:
            banner_h = max(0.8, 7.0 - banner_y)
        R(slide, 0.45, banner_y, 12.4, banner_h, PALETTE_LIGHT, rounded=True)
        T(slide, 0.7, banner_y + 0.1, 11.9, 0.4, '核心假说与实验设计', sz=Pt(18), bold=True, color=PALETTE_PRIMARY)
        blines = []
        if hypothesis:
            blines.append('假说: ' + hypothesis)
        if experiment:
            blines.append('实验: ' + experiment)
        M(slide, 0.7, banner_y + 0.55, 11.9, banner_h - 0.7, blines, sz=BODY_SIZE, color=DARK)

    return slide


def build_summary_slide(prs, data):
    """
    构建"结果总结"页。
    布局：顶部横向流程箭头图 + 中部证据卡片 + 底部结论横幅。

    data 结构:
      title: 幻灯片标题
      flow_steps: [{text, color}, ...]  — 流程步骤
      evidence_cards: [{title, detail}, ...]  — 证据卡片
      conclusions: [字符串列表]  — 结论要点
    """
    slide = prs.slides.add_slide(blank_layout(prs))
    white_bg(slide)
    title_bar(slide, data.get('title', '研究总结'))

    # ── 顶部流程步骤（深色卡片，高度自适应文本）──
    flow_steps = data.get('flow_steps', [])
    n = len(flow_steps)
    flow_end_y = 1.15  # 无流程步骤时的默认值
    if n > 0:
        max_cols = 5
        n_rows = (n + max_cols - 1) // max_cols
        per_row = min(n, max_cols)
        base_y = 1.15
        fs_gap = 0.2  # 行间距
        for row in range(n_rows):
            row_steps = flow_steps[row * per_row : (row + 1) * per_row]
            rn = len(row_steps)
            bw = max(1.2, (12.5 - (rn - 1) * 0.08) / rn)
            bw = min(bw, 10.0 if rn == 1 else 2.2)  # 单步时拉宽到 10in，多步时保持 2.2in
            x0 = (13.33 - (rn * (bw + 0.08) - 0.08)) / 2
            # 计算本行最大卡片高度（按文本折行数 + 上下padding 0.15in）
            fs_h = 0.8  # 最小 0.8in
            fs_text_w = bw - 0.16  # 文本区宽度
            fs_font = Pt(10)
            for step in row_steps:
                n_lines = _est_wrapped_lines([step['text']], fs_font, fs_text_w)
                need_h = n_lines * _line_h(fs_font) + 0.2
                fs_h = max(fs_h, min(need_h, 2.5))  # 上限 2.5in
            for i, step in enumerate(row_steps):
                x = x0 + i * (bw + 0.08)
                color = parse_color(step.get('color', 'teal'))
                R(slide, x, base_y, bw, fs_h, color, rounded=True)
                # 截断过长文本
                step_text = truncate_lines([step['text']], fs_font, fs_text_w, fs_h - 0.2)[0]
                T(slide, x + 0.08, base_y + 0.1, fs_text_w, fs_h - 0.2, step_text,
                   sz=fs_font, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
                if i < rn - 1:
                    mid_y = base_y + fs_h / 2 - 0.15
                    T(slide, x + bw + 0.005, mid_y, 0.1, 0.3, '→', sz=Pt(16), bold=True,
                       color=PALETTE_PRIMARY, align=PP_ALIGN.CENTER)
            base_y += fs_h + fs_gap
        flow_end_y = base_y - fs_gap  # 流程步骤结束 Y

    # ── 中部证据卡片（位置紧跟流程步骤）──
    evidence_cards = data.get('evidence_cards', [])
    ev_end_y = flow_end_y + 0.2
    if evidence_cards:
        nc = len(evidence_cards)
        cw = min(3.95, (12.5 - (nc - 1) * 0.1) / nc)
        detail_w = cw - 0.4
        max_detail_wrapped = 1
        for card in evidence_cards:
            detail_text = card.get('detail', '')
            if detail_text.strip():
                n = _est_wrapped_lines([detail_text], BODY_SIZE, detail_w)
                max_detail_wrapped = max(max_detail_wrapped, n)
        body_line_h = _line_h(BODY_SIZE)
        ev_h = 0.55 + (max_detail_wrapped + 1) * body_line_h + 0.1
        ev_h = max(1.6, ev_h)
        ev_h = min(4.5, ev_h)
        ev_base_y = flow_end_y + 0.2
        for i, card in enumerate(evidence_cards):
            cx = 0.45 + i * (cw + 0.1)
            R(slide, cx, ev_base_y, cw, ev_h, PALETTE_LIGHT, rounded=True)
            T(slide, cx + 0.2, ev_base_y + 0.1, cw - 0.4, 0.35, card['title'],
               sz=BODY_SIZE, bold=True, color=PALETTE_PRIMARY)
            detail_text = card.get('detail', '')
            detail_trunc = truncate_lines([detail_text], BODY_SIZE, detail_w, ev_h - 0.6)[0]
            T(slide, cx + 0.2, ev_base_y + 0.5, detail_w, ev_h - 0.6, detail_trunc,
               sz=BODY_SIZE, color=DARK)
        ev_end_y = ev_base_y + ev_h + 0.15

    # ── 底部结论横幅（位置紧跟证据卡片）──
    conclusions = data.get('conclusions', [])
    if conclusions:
        conc_h = max(1.5, 0.6 + len(conclusions) * 0.45)
        conc_h = min(3.5, conc_h)
        conc_y = ev_end_y
        if conc_y + conc_h > 7.0:
            conc_h = max(0.8, 7.0 - conc_y)
        R(slide, 0.45, conc_y, 12.4, conc_h, PALETTE_LIGHT, rounded=True)
        T(slide, 0.7, conc_y + 0.1, 11.9, 0.35, '核心结论', sz=Pt(18), bold=True, color=PALETTE_PRIMARY)
        conc_trunc = truncate_lines(conclusions, BODY_SIZE, 11.9, conc_h - 0.6)
        M(slide, 0.7, conc_y + 0.55, 11.9, conc_h - 0.6, conc_trunc,
           sz=BODY_SIZE, color=DARK)

    return slide


def build_discussion1_slide(prs, data):
    """
    构建"讨论1"页 —— 编号圆形列表布局。
    每个条目：左侧青色圆形编号 + 右侧标题 + 下方详细描述。
    条目之间用细灰色分隔线隔开。

    data 结构:
      title: 幻灯片标题
      items: [{number, title, detail}, ...]  — 讨论条目列表
    """
    slide = prs.slides.add_slide(blank_layout(prs))
    white_bg(slide)
    title_bar(slide, data.get('title', '4 讨论'))

    items = data.get('items', [])
    # 根据条目数动态调整间距，避免文字重叠
    n_items = max(len(items), 1)
    available_h = 5.8  # title_bar 下方可用高度
    row_h = min(1.48, available_h / n_items)  # 每行高度，最多 1.48in
    for i, item in enumerate(items):
        y = 1.15 + i * row_h
        detail_h = max(0.6, row_h - 0.55)  # 详细描述高度自适应
        # 圆形编号
        c = slide.shapes.add_shape(9, Inches(0.55), Inches(y + 0.05), Inches(0.45), Inches(0.45))
        c.fill.solid(); c.fill.fore_color.rgb = PALETTE_PRIMARY; c.line.fill.background()
        T(slide, 0.55, y + 0.08, 0.45, 0.4, item.get('number', str(i + 1)),
           sz=Pt(14), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        # 标题（加宽到 6.0in 避免中文标题被截断）
        T(slide, 1.25, y, 6.0, 0.35, item['title'], sz=BODY_SIZE, bold=True, color=PALETTE_PRIMARY)
        # 详细描述（截断过长文本）
        detail_trunc = truncate_lines([item['detail']], BODY_SIZE, 11.5, detail_h)[0]
        T(slide, 1.25, y + 0.35, 11.5, detail_h, detail_trunc, sz=BODY_SIZE, color=DARK)
        # 分隔线
        if i < len(items) - 1:
            sep_y = y + row_h - 0.08
            sep = slide.shapes.add_shape(1, Inches(1.25), Inches(sep_y), Inches(11.0), Pt(0.5))
            sep.fill.solid(); sep.fill.fore_color.rgb = RGBColor(0xDD, 0xDD, 0xDD)
            sep.line.fill.background()

    return slide


def build_discussion2_slide(prs, data):
    """
    构建"讨论2"页 —— 左右双栏布局。
    左侧：局限性列表（浅青色背景）
    右侧：临床意义/未来方向（浅绿色背景）

    data 结构:
      title: 幻灯片标题
      left_title: 左栏标题
      left_items: 左栏条目列表
      right_title: 右栏标题
      right_items: 右栏条目列表
    """
    slide = prs.slides.add_slide(blank_layout(prs))
    white_bg(slide)
    title_bar(slide, data.get('title', '4 局限性与临床意义'))

    left_title = data.get('left_title', '局限性')
    left_items = data.get('left_items', [])
    right_title = data.get('right_title', '临床意义')
    right_items = data.get('right_items', [])

    # 卡片高度拉满：title_bar 下方 (~1.0in) 到页脚上方 (~6.8in)，不截断
    n_left = max(len(left_items), 1); n_right = max(len(right_items), 1)
    max_items = max(n_left, n_right, 1)
    col_h = min(5.7, max(2.8, max_items * 1.0 + 0.6))
    text_area_h = col_h - 0.7  # 标题占用 0.4 + 间距 0.3

    # 左栏
    R(slide, 0.4, 1.1, 6.0, col_h, PALETTE_LIGHT, rounded=True)
    T(slide, 0.7, 1.25, 5.4, 0.4, left_title, sz=Pt(20), bold=True, color=PALETTE_PRIMARY)
    M(slide, 0.7, 1.8, 5.4, text_area_h, left_items, sz=BODY_SIZE, color=DARK)

    # 右栏
    R(slide, 6.8, 1.1, 6.1, col_h, PALETTE_WARM, rounded=True)
    T(slide, 7.1, 1.25, 5.5, 0.4, right_title, sz=Pt(20), bold=True, color=PALETTE_ACCENT2)
    M(slide, 7.1, 1.8, 5.5, text_area_h, right_items, sz=BODY_SIZE, color=DARK)

    return slide


def build_paper_info_slide(prs, data):
    """
    论文信息页：左侧 PDF 首页截图 + 右侧原文链接（自动搜索）。
    data: {pdf_path, paper_title, extra_text}
    """
    import urllib.request, urllib.parse
    slide = prs.slides.add_slide(blank_layout(prs))
    white_bg(slide)
    title_bar(slide, '原文信息')

    pdf_path = data.get('pdf_path', '')
    pm = data.get('paper_meta', {})
    paper_title = pm.get('title_en') or data.get('paper_title', '')

    # ── 左侧：PDF 首页渲染 ──
    first_page_img = None
    if pdf_path and os.path.exists(pdf_path):
        try:
            import pypdfium2 as pdfium
            pdf = pdfium.PdfDocument(pdf_path)
            bitmap = pdf[0].render(scale=1.5)
            img_path = pdf_path.replace('.pdf', '_page1.jpg')
            bitmap.to_pil().save(img_path, 'JPEG', quality=85)
            first_page_img = img_path
        except Exception as e:
            print(f"  WARNING: PDF render: {e}")

    if first_page_img:
        P(slide, first_page_img, 0.3, 1.1, 5.5, max_h=5.8)
    else:
        T(slide, 0.5, 2.5, 5.5, 1.0, '(PDF 首页渲染失败)', sz=Pt(12), color=RGBColor(0x99, 0x99, 0x99))

    # ── 右侧：带标签的元数据表格 ──
    right_x = 6.5
    right_w = 6.3
    label_w = 0.85   # 标签列宽
    value_x = right_x + label_w + 0.1  # 内容起始 x
    value_w = right_w - label_w - 0.1  # 内容列宽
    label_sz = Pt(11)     # 标签字号
    value_sz = Pt(11)     # 内容字号
    label_color = PALETTE_PRIMARY
    value_color = DARK
    row_h = 0.28  # 行高
    cur_y = 1.15

    def _row(label, value, lines=1):
        """画一行：标签（加粗主题色）+ 值。返回实际占用的 Y 增量。"""
        nonlocal cur_y
        h = max(0.25, lines * row_h + 0.05)
        T(slide, right_x, cur_y, label_w, h, label, sz=label_sz, bold=True, color=label_color,
           align=PP_ALIGN.RIGHT)
        T(slide, value_x, cur_y, value_w, h, value, sz=value_sz, color=value_color)
        cur_y += h + 0.02

    # 标题
    if paper_title:
        T(slide, right_x, cur_y, right_w, 0.55, paper_title, sz=Pt(16), bold=True, color=PALETTE_PRIMARY)
        cur_y += 0.55

    # ── 表格行 ──
    authors = pm.get('authors', '')
    if authors:
        _row('作者', authors)

    journal = pm.get('citation', pm.get('journal', ''))
    if journal:
        _row('期刊', journal)

    year = pm.get('year', '')
    if year:
        _row('年份', str(year))

    doi = pm.get('doi', '')
    if doi:
        _row('DOI', doi)

    keywords = pm.get('keywords', '')
    if keywords:
        _row('关键词', keywords)

    # 分隔线
    cur_y += 0.05
    sep = slide.shapes.add_shape(1, Inches(right_x), Inches(cur_y),
                                  Inches(right_w), Pt(0.5))
    sep.fill.solid(); sep.fill.fore_color.rgb = RGBColor(0xCC, 0xCC, 0xCC); sep.line.fill.background()
    cur_y += 0.2

    # 摘要
    abstract = pm.get('abstract', data.get('extra_text', ''))
    if abstract:
        T(slide, right_x, cur_y, label_w, 0.3, '摘要', sz=label_sz, bold=True, color=label_color,
           align=PP_ALIGN.RIGHT)
        max_abs_h = 6.8 - cur_y - 0.35
        if len(abstract) > 500:
            abstract = abstract[:497] + '...'
        abs_text = truncate_lines([abstract], Pt(11), value_w, max_abs_h)[0] if max_abs_h > 0.5 else abstract[:200]
        T(slide, value_x, cur_y + 0.3, value_w, max_abs_h, abs_text, sz=Pt(11), color=value_color)

    # 无元数据时的兜底：原文链接
    if not pm and paper_title:
        link_lines = []
        try:
            q = urllib.parse.quote(paper_title[:200])
            u = f'https://api.semanticscholar.org/graph/v1/paper/search?query={q}&limit=1'
            req = urllib.request.Request(u, headers={'User-Agent': 'PaperDeck/1.0'})
            with urllib.request.urlopen(req, timeout=10) as resp:
                import json as _json
                d = _json.loads(resp.read())
                if d.get('data'):
                    p = d['data'][0]
                    pid = p.get('paperId', '')
                    eid = p.get('externalIds', {})
                    link_lines.append(f'Semantic Scholar: https://api.semanticscholar.org/CorpusID:{pid}')
                    if eid.get('DOI'):
                        link_lines.append(f'DOI: https://doi.org/{eid["DOI"]}')
                    if eid.get('ArXiv'):
                        link_lines.append(f'arXiv: https://arxiv.org/abs/{eid["ArXiv"]}')
        except Exception:
            pass
        if not link_lines:
            q = urllib.parse.quote(paper_title[:200])
            link_lines.append(f'Google Scholar: https://scholar.google.com/scholar?q={q}')
        T(slide, right_x, cur_y, right_w, 0.4, '原文链接', sz=Pt(14), bold=True, color=PALETTE_PRIMARY)
        M(slide, right_x, cur_y + 0.5, right_w, 3.0, link_lines, sz=Pt(12), color=DARK)

    return slide
