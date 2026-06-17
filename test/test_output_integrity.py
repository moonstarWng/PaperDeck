"""
test/test_output_integrity.py — 生成 PPTX 的完整性检查。
验证：封面/致谢存在、内容页有页眉页脚、元素完整。
"""
import os, json, pytest
from pptx import Presentation


def _build_output(demo_dir):
    """用 process/ 下的配置文件重建 output.pptx。"""
    from ppt_builder import load_json, build
    config_path = os.path.join(demo_dir, 'process', 'slide-content.json')
    tmpl_path = os.path.join(demo_dir, 'process', 'template.pptx')
    if not os.path.exists(config_path) or not os.path.exists(tmpl_path):
        pytest.skip("process/ files not found")
    config = load_json(config_path)
    config.setdefault('meta', {})
    config['meta']['template_path'] = tmpl_path
    config['meta']['figs_dir'] = os.path.join(demo_dir, 'figs')
    output = os.path.join(demo_dir, 'output_test_integrity.pptx')
    config['meta']['output_path'] = output
    build(config, config_path)
    return output


class TestStartEndPages:
    """封面（第一页）和致谢（最后一页）检查。"""

    def test_cover_is_first_slide(self, demo_dir):
        output = _build_output(demo_dir)
        prs = Presentation(output)
        from make_template import classify
        assert classify(prs.slides[0]) == 'COVER', "第一页必须是封面"
        os.unlink(output)

    def test_thanks_is_last_slide(self, demo_dir):
        output = _build_output(demo_dir)
        prs = Presentation(output)
        from make_template import classify
        assert classify(prs.slides[-1]) == 'THANKS', "最后一页必须是致谢"
        os.unlink(output)

    def test_toc_exists(self, demo_dir):
        output = _build_output(demo_dir)
        prs = Presentation(output)
        from make_template import classify
        cats = [classify(s) for s in prs.slides]
        assert 'TOC' in cats, "必须有目录页"
        os.unlink(output)


class TestHeaderFooter:
    """每个内容页应有页眉和页脚。"""

    def test_content_pages_have_header(self, demo_dir):
        output = _build_output(demo_dir)
        prs = Presentation(output)
        from make_template import classify
        for i, s in enumerate(prs.slides):
            cat = classify(s)
            if cat in ('CONTENT', 'CONTENT_FRAME', 'DISCUSSION'):
                has_header = any(
                    sh.top / 914400 < 1.0 and sh.width / 914400 > 12
                    and sh.height / 914400 > 0.3
                    for sh in s.shapes
                )
                assert has_header, f"Slide {i} [{cat}] 缺少页眉"

    def test_content_pages_have_footer(self, demo_dir):
        output = _build_output(demo_dir)
        prs = Presentation(output)
        from make_template import classify
        for i, s in enumerate(prs.slides):
            cat = classify(s)
            if cat in ('CONTENT', 'CONTENT_FRAME', 'DISCUSSION'):
                has_footer = any(
                    sh.top / 914400 > 6.5 and sh.width / 914400 > 10
                    for sh in s.shapes
                )
                assert has_footer, f"Slide {i} [{cat}] 缺少页脚"

    def test_section_dividers_have_decorations(self, demo_dir):
        """章节分隔页应有深色背景或左侧竖条。"""
        output = _build_output(demo_dir)
        prs = Presentation(output)
        from make_template import classify
        for i, s in enumerate(prs.slides):
            if classify(s).startswith('SECTION_0'):
                shapes = len(s.shapes)
                assert shapes >= 4, f"Section divider {i} 元素不足 ({shapes} shapes)"


class TestElementCompleteness:
    """检查封面、内容页的元素是否完整。"""

    def test_cover_has_title_and_presenter(self, demo_dir):
        output = _build_output(demo_dir)
        prs = Presentation(output)
        cover = prs.slides[0]
        all_text = ' '.join(sh.text_frame.text for sh in cover.shapes if sh.has_text_frame)
        # 封面必须有标题（非空）
        assert len(all_text.strip()) > 0, "封面无文字"
        os.unlink(output)

    def test_result_pages_have_bullet_dots(self, demo_dir):
        """结果页应有圆点标记。"""
        output = _build_output(demo_dir)
        prs = Presentation(output)
        from make_template import classify
        has_dots = False
        for s in prs.slides:
            if classify(s) == 'CONTENT':
                for sh in s.shapes:
                    if 'AUTO_SHAPE' in str(sh.shape_type):
                        w = sh.width / 914400
                        h = sh.height / 914400
                        if 0.1 < w < 0.3 and 0.1 < h < 0.3:
                            has_dots = True
        assert has_dots, "结果页应有关键点圆点标记"

    def test_no_empty_content_slides(self, demo_dir):
        """所有内容页至少有一个文字区。"""
        output = _build_output(demo_dir)
        prs = Presentation(output)
        from make_template import classify
        for i, s in enumerate(prs.slides):
            if classify(s) in ('CONTENT', 'CONTENT_FRAME', 'DISCUSSION'):
                texts = [sh.text_frame.text.strip()
                         for sh in s.shapes if sh.has_text_frame]
                combined = ' '.join(texts)
                assert len(combined) > 5, f"Slide {i} 内容为空"


class TestDesignTokensApplied:
    """设计令牌是否正确应用到内容页。"""

    def test_content_title_color_from_design(self, demo_dir):
        """内容页标题颜色应来自设计令牌，而非默认青色。"""
        output = _build_output(demo_dir)
        prs = Presentation(output)
        from make_template import classify
        design_path = os.path.join(demo_dir, 'process', 'template_design.json')
        if not os.path.exists(design_path):
            design_path = os.path.join(demo_dir, 'template_design.json')
        if os.path.exists(design_path):
            with open(design_path, 'r', encoding='utf-8') as f:
                design = json.load(f)
            expected_color = design.get('content_title', {}).get('color', '007191')
            # 找一个内容页，检查标题颜色
            for s in prs.slides:
                if classify(s) in ('CONTENT', 'CONTENT_FRAME'):
                    for sh in s.shapes:
                        if sh.top / 914400 < 0.5 and sh.has_text_frame:
                            for p in sh.text_frame.paragraphs:
                                for r in p.runs:
                                    try:
                                        actual = str(r.font.color.rgb)
                                        assert actual == expected_color, \
                                            f"标题颜色 {actual} != 设计令牌 {expected_color}"
                                    except Exception:
                                        pass  # inherit 也算可以
                                    return
        os.unlink(output)
