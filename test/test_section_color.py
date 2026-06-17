"""
test/test_section_color.py — 章节页 bg_color 提取：确保永远取最暗的背景形状。
"""
import os, pytest
from pptx import Presentation
from pptx.util import Inches
from pptx.dml.color import RGBColor


class TestDetectSectionColors:
    """验证 _detect_section_colors 的 bg_color 始终是深色。"""

    def test_returns_dark_bg_on_real_template(self, demo_dir):
        """在 process/template.pptx 章节页上，提取的 bg_color 必须是深色（亮度 < 0.3）。"""
        from ppt_builder import _detect_section_colors
        from ppt_layout import _luminance
        tmpl = os.path.join(demo_dir, 'process', 'template.pptx')
        if not os.path.exists(tmpl):
            pytest.skip("process/template.pptx not found")
        prs = Presentation(tmpl)
        from make_template import classify
        found = None
        for slide in prs.slides:
            if classify(slide).startswith('SECTION_'):
                found = slide
                break
        if found is None:
            pytest.skip("模板无章节页")
        bg_color, _ = _detect_section_colors(found)
        lum = _luminance(bg_color)
        assert lum < 0.3, f"bg_color 亮度 {lum:.3f} 过高 (应为深色)"

    def test_returns_dark_bg_on_demo_template(self, demo_dir):
        """在 demo_paper_汇报.pptx 章节页上，bg_color 必须是深色。"""
        from ppt_builder import _detect_section_colors
        from ppt_layout import _luminance
        tmpl = os.path.join(demo_dir, 'demo_paper_汇报.pptx')
        if not os.path.exists(tmpl):
            pytest.skip("demo_paper_汇报.pptx not found")
        prs = Presentation(tmpl)
        from make_template import classify
        found = None
        for slide in prs.slides:
            if classify(slide).startswith('SECTION_'):
                found = slide
                break
        if found is None:
            pytest.skip("模板无章节页")
        bg_color, _ = _detect_section_colors(found)
        lum = _luminance(bg_color)
        assert lum < 0.3, f"bg_color 亮度 {lum:.3f} 过高 (应为深色)"

    def test_picks_darkest_when_multiple_full_width_shapes(self):
        """多个全幅形状时，取最暗的，不取第一个。"""
        from ppt_builder import _detect_section_colors
        from ppt_layout import _luminance
        # 构造一个带两个全幅形状的幻灯片：先白后黑
        prs = Presentation()
        ly = prs.slide_layouts[0]
        slide = prs.slides.add_slide(ly)
        # 先加白色全幅形状（模拟浅色遮罩）
        white_shape = slide.shapes.add_shape(
            1, Inches(0), Inches(0), Inches(13.33), Inches(7.5))
        white_shape.fill.solid()
        white_shape.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        # 再加黑色全幅形状（真正的背景）
        black_shape = slide.shapes.add_shape(
            1, Inches(0), Inches(0), Inches(13.33), Inches(7.5))
        black_shape.fill.solid()
        black_shape.fill.fore_color.rgb = RGBColor(0x1A, 0x1A, 0x1A)
        bg_color, _ = _detect_section_colors(slide)
        lum = _luminance(bg_color)
        assert lum < 0.3, f"多全幅形状应取最暗的，但亮度 {lum:.3f}（可能取了白的）"
