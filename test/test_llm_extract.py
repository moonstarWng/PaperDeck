"""
test/test_llm_extract.py — dump_slide_info / analyze_and_cache / load_template_cache / apply_placeholders
"""
import os, json, pytest
from pptx import Presentation


class TestDumpSlideInfo:
    def test_dump_all_slides(self, template_path):
        from llm_template_extract import dump_all_slides
        prs = Presentation(template_path)
        data = dump_all_slides(prs)
        assert len(data) == len(prs.slides)
        for sd in data:
            assert 'shapes' in sd
            assert 'slide_index' in sd
            assert len(sd['shapes']) > 0

    def test_dump_includes_text_and_font(self, template_path):
        from llm_template_extract import dump_all_slides
        from make_template import classify
        prs = Presentation(template_path)
        data = dump_all_slides(prs)
        # COVER 页应有文字形状
        for sd in data:
            if sd['slide_index'] == 0:
                has_text = any('text' in s for s in sd['shapes'])
                assert has_text


class TestCache:
    def test_analyze_and_cache_creates_file(self, template_path):
        from llm_template_extract import analyze_and_cache, load_template_cache
        cache = analyze_and_cache(template_path)
        assert cache is not None
        assert 'slides' in cache
        assert len(cache['slides']) >= 4
        for s in cache['slides']:
            assert 'classification' in s
            assert s['classification'] in ('COVER', 'TOC', 'SECTION', 'THANKS', 'CONTENT')

    def test_cache_hit(self, template_path):
        from llm_template_extract import load_template_cache
        cache = load_template_cache(template_path)
        assert cache is not None, "Cache should exist after analyze_and_cache"

    def test_cache_has_design_tokens(self, template_path):
        from llm_template_extract import load_template_cache
        cache = load_template_cache(template_path)
        assert 'design_tokens' in cache
        dt = cache['design_tokens']
        # 至少应有这些字段
        for key in ('footer_bars', 'content_title', 'section_number', 'toc_number'):
            assert key in dt, f"Missing design token: {key}"


class TestApplyPlaceholders:
    def test_title_becomes_placeholder(self, template_path, demo_dir):
        from llm_template_extract import apply_placeholders, dump_slide_info
        prs = Presentation(template_path)
        sd = dump_slide_info(prs.slides[0])
        # 模拟封面有标题
        emap = {}
        for i, s in enumerate(sd['shapes']):
            if s.get('text') and len(s.get('text', '')) > 10:
                emap[i] = 'title'
            else:
                emap[i] = 'decoration'
        apply_placeholders(prs, 0, emap)
        # 检查文字是否被替换
        for i, sh in enumerate(prs.slides[0].shapes):
            if i in emap and emap[i] == 'title' and sh.has_text_frame:
                assert '此处填充标题' in sh.text_frame.text or '论文标题' in sh.text_frame.text
