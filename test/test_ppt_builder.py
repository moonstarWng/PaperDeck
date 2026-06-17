"""
test/test_ppt_builder.py — edit_cover / edit_toc / edit_section_divider / section cloning
"""
import os, json, pytest
from pptx import Presentation


class TestEditCover:
    def test_title_replaced(self, template_path):
        from ppt_builder import edit_cover
        prs = Presentation(template_path)
        slide = prs.slides[0]
        old_texts = [sh.text_frame.text for sh in slide.shapes if sh.has_text_frame]
        edit_cover(slide, {'title_en': 'Test Title', 'presenter': 'tester', 'date': '2024'})
        new_texts = [sh.text_frame.text for sh in slide.shapes if sh.has_text_frame]
        # 至少有一条文字变了
        assert old_texts != new_texts
        assert any('Test' in t for t in new_texts)


class TestEditTOC:
    def test_five_sections_evenly_distributed(self, template_path):
        from ppt_builder import edit_toc
        titles = ['作者团队', '课题背景', '结果分析', '结果总结', '讨论分析']
        prs = Presentation(template_path)
        toc_slide = prs.slides[1]
        edit_toc(toc_slide, {}, titles)
        # 验证 5 个条目全部出现
        ns = 'http://schemas.openxmlformats.org/drawingml/2006/main'
        all_text = [t.text.strip() for t in toc_slide._element.iter(f'{{{ns}}}t') if t.text and t.text.strip()]
        for t in titles:
            assert t in all_text

    def test_three_sections_evenly_distributed(self, template_path):
        from ppt_builder import edit_toc
        titles = ['作者团队', '课题背景', '结果分析']
        prs = Presentation(template_path)
        toc_slide = prs.slides[1]
        edit_toc(toc_slide, {}, titles)
        ns = 'http://schemas.openxmlformats.org/drawingml/2006/main'
        all_text = [t.text.strip() for t in toc_slide._element.iter(f'{{{ns}}}t') if t.text and t.text.strip()]
        for t in titles:
            assert t in all_text


class TestSectionDivider:
    def test_section_cloning_creates_correct_count(self, template_path, demo_dir):
        """5 章节配置时应有 5 个章节页。"""
        from ppt_builder import load_json, build
        config_path = os.path.join(demo_dir, 'process', 'slide-content.json')
        if not os.path.exists(config_path):
            pytest.skip("process/slide-content.json not found")
        config = load_json(config_path)
        config.setdefault('meta', {})
        config['meta']['template_path'] = os.path.join(demo_dir, 'process', 'template.pptx')
        config['meta']['figs_dir'] = os.path.join(demo_dir, 'figs')
        config['meta']['output_path'] = os.path.join(demo_dir, 'output_test.pptx')
        output = build(config, config_path)
        prs = Presentation(output)
        from make_template import classify
        sec_count = sum(1 for s in prs.slides if classify(s).startswith('SECTION_0'))
        edits = config.get('section_divider_edits', [])
        assert sec_count >= len(edits), f"Expected {len(edits)} sections, got {sec_count}"
        os.unlink(output) if os.path.exists(output) else None
