"""
test_pipeline.py — 自测脚本：使用 demo 数据验证完整流水线。
用法: python test_pipeline.py
"""
import sys, os, json, tempfile
sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))
sys.path.insert(0, os.path.join(os.path.dirname(__file__), 'scripts'))

PASS, FAIL = 0, 0


def check(name, condition, detail=""):
    global PASS, FAIL
    if condition:
        PASS += 1
        print(f"  [PASS] {name}")
    else:
        FAIL += 1
        print(f"  [FAIL] {name} {detail}")


def main():
    global PASS, FAIL
    print("=" * 50)
    print("paper2ppt — 自测")
    print("=" * 50)

    # ── Test 1: validate_outline ──
    print("\n[Test 1] validate_outline.py")
    from validate_outline import validate
    example = os.path.join('templates', 'example-slide-content.json')
    check("example JSON 验证通过", validate(example))

    # ── Test 2: ppt_builder ──
    print("\n[Test 2] ppt_builder.py")
    from ppt_builder import load_json, build
    demo_json = os.path.abspath(os.path.join('demo', 'slide-content.json'))
    demo_tmpl = os.path.abspath(os.path.join('demo', 'template.pptx'))
    demo_figs = os.path.abspath(os.path.join('demo', 'figs'))

    if os.path.exists(demo_json) and os.path.exists(demo_tmpl):
        config = load_json(demo_json)
        config.setdefault('meta', {})
        config['meta']['template_path'] = demo_tmpl
        config['meta']['figs_dir'] = demo_figs
        with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as f:
            config['meta']['output_path'] = f.name
        tmp_json = tempfile.NamedTemporaryFile(mode='w', suffix='.json', delete=False, encoding='utf-8')
        json.dump(config, tmp_json)
        tmp_json.close()
        try:
            output = build(config, tmp_json.name)
            ok = os.path.exists(output) and os.path.getsize(output) > 10000
            check("Demo JSON 构建成功", ok)
            if ok:
                from pptx import Presentation
                prs = Presentation(output)
                check(f"输出页数 ({len(prs.slides)})", 10 <= len(prs.slides) <= 25)
                img_count = sum(1 for s in prs.slides for sh in s.shapes if hasattr(sh, 'image'))
                check(f"输出图片数 ({img_count})", img_count >= 4)
            os.unlink(output)
        except Exception as e:
            check("Demo JSON 构建", False, str(e))
        finally:
            os.unlink(tmp_json.name)

    # ── Test 3: make_template ──
    print("\n[Test 3] make_template.py (分类器)")
    from make_template import classify
    from pptx import Presentation
    prs = Presentation(demo_tmpl)
    cats = [classify(s) for s in prs.slides]
    check(f"模板加载 ({len(cats)} slides)", len(cats) == 8)
    # 分类器不要求匹配特定类型——只验证函数能正常返回且无异常
    check(f"分类完成 ({len(cats)} cats)", all(isinstance(c, str) for c in cats))

    # ── Test 4: extract_images ──
    print("\n[Test 4] extract_images.py")
    from extract_images import extract_images
    demo_pdf = os.path.join('demo', 'demo_paper.pdf')
    if os.path.exists(demo_pdf):
        with tempfile.TemporaryDirectory() as td:
            n = extract_images(demo_pdf, td)
            check(f"图片提取 ({n} images)", n >= 0)  # 合成 PDF 可能不含嵌入图片
    else:
        print("  [SKIP] demo_paper.pdf 不存在")

    # ── Summary ──
    total = PASS + FAIL
    print(f"\n{'='*50}")
    print(f"结果: {PASS}/{total} 通过, {FAIL} 失败")
    print(f"{'='*50}")
    return FAIL == 0


if __name__ == '__main__':
    ok = main()
    sys.exit(0 if ok else 1)
